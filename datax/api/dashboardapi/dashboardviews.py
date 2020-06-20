from account.models import user_tag

import shutil
from pathlib import  Path
import base64
import io
import json, re
import random
import time
from datetime import date
import urllib

import xlsxwriter
from PIL import Image, ImageDraw, ImageFont
from bs4 import BeautifulSoup
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as pptInches
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from rest_framework import permissions
from rest_framework import request
from rest_framework.decorators import api_view, permission_classes
from rest_framework.response import Response
from django.contrib.auth import get_user
import pandas as pd

from account.models import user_tag
from api import utils
from api.mail_utils import *
from common import tools
from common.head import LIMIT
from connect.models import olap, source, olapcolumn, dispatchlog, monitor, monitorDetail, olapextcols
from common import globalvariable as glvr
from common.constantcode import DBTypeCode
from connect.sqltool import stRestoreBySourceId

from dashboard.models import *
from dashboard.models import data_dictionary
import sys
from connect.models import maillogs, systemmessage
from account.models import sys_userextension
import os, zipfile, tempfile
from django.http import StreamingHttpResponse
from channels import Channel, Group
import logging
from common.constantcode import DBTypeCode,LoggerCode
import pandas as pd
import datetime
logger = logging.getLogger(LoggerCode.DJANGOINFO.value)

# from apscheduler.scheduler import Scheduler


# 自定义sql查询结束
#
# class UserViewSet(viewsets.ModelViewSet):
#     """
#     API endpoint that allows users to be viewed or edited.
#     """
#     queryset = User.objects.all().order_by('-date_joined')
#     serializer_class = UserSerializer
#
#
# class GroupViewSet(viewsets.ModelViewSet):
#     """
#     API endpoint that allows groups to be viewed or edited.
#     """
#     queryset = Group.objects.all()
#     serializer_class = GroupSerializer
@api_view(http_method_names=['POST'])
def deletedatapermtag(request):
    id = request.data.get('id')
    resultObj = tools.successMes()
    if id:
        try:
            utag = user_tag.objects.get(id=id)
            utag.delete()
        except Exception as e:
            resultObj = tools.errorMes(e.args)
    else:
        resultObj = tools.errorMes('id为空')
    return Response(resultObj)


@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
# 检查数据标签是否重复
def checktagname(request):
    tagname = request.data.get('name')
    resultObj = tools.successMes()
    try:
        dtag = user_tag.objects.filter(name=tagname).first()
        if not tagname:
            resultObj['code'] = 1
        if dtag:
            resultObj['code'] = 0
            resultObj['msg'] = '标签名重复'
        else:
            resultObj['code'] = 1
    except Exception as e:
        resultObj = tools.errorMes(e.args)
        resultObj['code'] = 1
    return Response(resultObj)


@api_view(http_method_names=['POST'])
def datapermtagadd(request):
    id = request.data.get("id")
    tagname = request.data.get("tagname")
    tagremark = request.data.get("tagremark")
    resultObj = tools.successMes()
    # 编辑
    if id:
        dtag = user_tag.objects.get(id=id)
        dtag.name = tagname
        dtag.remark = tagremark
        try:
            dtag.save()
            resultObj["code"] = 1
        except Exception as e:
            resultObj = tools.errorMes(e.args)
            resultObj["code"] = 0
    # 新增
    else:
        try:
            user_tag.objects.get_or_create(name=tagname, remark=tagremark)
            resultObj["code"] = 1
        except Exception as e:
            resultObj = tools.errorMes(e.args)
            resultObj["code"] = 0
    return Response(resultObj)


@api_view(http_method_names=['GET'])
def datapermtagview(request):
    rs = {}
    if 'id' in request.GET:
        id = request.GET['id']
        datatag = user_tag.objects.get(id=id)
        rs['id'] = datatag.id
        rs['tagname'] = datatag.name
        rs['tagremark'] = datatag.remark
    return Response({'data': rs})


# 获取单个
@api_view(http_method_names=['GET'])
def chartsbgconfigview(request):
    rs = {}
    if 'id' in request.GET:
        id = request.GET['id']
        chartsbg = chartsbgconfig.objects.get(id=id)
        rs['id'] = chartsbg.id
        rs['confname'] = chartsbg.confname
        rs['fileurl'] = chartsbg.fileurl
        rs['remark'] = chartsbg.remark
        rs['status'] = chartsbg.status
        rs['filecontent'] = chartsbg.filecontent
    return Response({'data': rs})

# 获取单个
@api_view(http_method_names=['GET'])
def scenebgconfigview(request):
    resultObj = tools.successMes()
    try:
        if request.GET['id'] == '0':
            resultSceneConfigObj = {}
            resultSceneConfigObj['id'] = ''
            resultSceneConfigObj['confname'] = ''
            resultSceneConfigObj['status'] = '0'
            resultSceneConfigObj['remark'] = ''
            resultSceneConfigObj['options'] = {"componentMargin":{"marginLeft":0,"marginTop":0},"globalStyle":{"color":"rgba(51, 51, 51, 1)","backgroundColor":"rgba(255,255,255, 1)","backgroundSize":"cover","backgroundRepeat":"no-repeat"},"toolBarStyle":{"default":True,"blue":False,"black":False},"filterStyle":{"backgroundColor":"rgba(255,255,255,1)","borderDisplay":False,"borderRadius":"0px","backgroundSize":"cover","backgroundRepeat":"no-repeat","currCharsBgPicSts":False},"chartStyle":{"backgroundColor":"rgba(255,255,255,1)","borderDisplay":False,"borderRadius":"0px","backgroundSize":"cover","backgroundRepeat":"no-repeat"},"syncFontColor":True,"commentsStyle":{"total":{"backgroundColor":"rgba(255, 255, 255, 1)"},"commentsSection":{"borderBottomColor":"rgba(210, 210, 210, 1)"},"name":{"color":"rgba(210, 179, 44, 1)"},"time":{"color":"rgba(153, 153, 153, 1)"},"text":{"color":"rgba(0, 0, 0, 1)","backgroundColor":"rgba(255, 255, 255, 1)","borderBottomColor":"rgba(210, 210, 210, 1)"},"editBtn":{"color":"rgba(210, 179, 44, 1)"},"replayBtn":{"color":"rgba(210, 179, 44, 1)"},"link":{"color":"rgba(67, 161, 248, 1)"},"more":{"color":"rgba(153, 153, 153, 1)","borderBottomColor":"rgba(210, 210, 210, 1)","height":"25px","cursor":"pointer"}}}
            resultSceneConfigObj['basicconfig'] = {"customConfig":{"customCode":{"code":""},"autoSave":{"active":False,"interval":60}}}
            resultObj['data'] = resultSceneConfigObj
            return Response(resultObj)
    except Exception as e:
        print(e)
    try:
        if 'id' in request.GET and request.GET['id'].strip():
            id = request.GET['id']
            sceneConfig = chartsbgconfig.objects.get(id=id)
            resultSceneConfigObj={}
            resultSceneConfigObj['id']=sceneConfig.id
            resultSceneConfigObj['confname']=sceneConfig.confname
            resultSceneConfigObj['status']=sceneConfig.status
            resultSceneConfigObj['remark']=sceneConfig.remark
            resultSceneConfigObj['options']=sceneConfig.configcontent
            resultSceneConfigObj['basicconfig']=sceneConfig.customconfig
            resultObj['data'] = resultSceneConfigObj
    except Exception as getSceneError:
        resultObj = tools.errorMes()
    return Response(resultObj)

# 获取chartsbgconfig，list
@api_view(http_method_names=['GET'])
def getChartsbgConfigList(request):
    where = ''
    if 'search' in request.GET:
        where = where + """ where confname like '%""" + request.GET['search'] + """%' """
    if 'status' in request.GET:
        if 'where' in where:
            where = where + """ and status = 0 """
        else:
            where = """ where status = 0 """
    if 'page' in request.GET:
        page = request.GET['page']
        offset = (int(request.GET['page']) - 1) * LIMIT
    session = utils.SqlUtils()
    allcnt = 0
    fetchall = []
    try:
        allcntQueryResult = session.getArrResult('select count(*) from dashboard_chartsbgconfig ' + where)
        allcnt = allcntQueryResult[0][0] if allcntQueryResult else 0
        if 'page' in request.GET:#如果有分页
            fetchall = session.getArrResult('select "id",confname,remark,status,creater,createtime,modifytime,configcontent,customconfig from dashboard_chartsbgconfig' +
                where + ' order by createtime desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset))
        else :
            fetchall = session.getArrResult('select "id",confname,remark,status,creater,createtime,modifytime,configcontent,customconfig from dashboard_chartsbgconfig' +
                where + ' order by createtime desc')
    except Exception as e:
        logger.error('---error---file:dashboardviews.py;method:getChartsbgConfigList;error=%s' % e)
    alldatas = []
    for obj in fetchall:
        singledt = {}
        singledt['id'] = obj[0]
        singledt['confname'] = obj[1]
        singledt['remark'] = obj[2]
        singledt['status'] = obj[3]
        singledt['creater'] = obj[4]
        if obj[5]:
            singledt['createtime'] = obj[5].strftime("%Y-%m-%d %H:%M:%S")
        else:
            singledt['createtime'] = obj[5]
        if obj[6]:
            singledt['modifytime'] = obj[6].strftime("%Y-%m-%d %H:%M:%S")
        else:
            singledt['modifytime'] = obj[6]
        singledt['options'] = obj[7]
        singledt['basicconfig'] = obj[8]
        alldatas.append(singledt)
    return Response({'total': allcnt, 'rows': alldatas})

# 保存
@api_view(http_method_names=['POST'])
def savescenebackgroundconfig(request):  # 接收用户传入的数据，将css内容写入新的文件保存到制定目录下，将必要数据持久化到数据库
    resultObj = tools.successMes()
    try :
        id = '0';  ##新增
        if 'id' in request.data and request.data['id'].strip():
            id = request.data['id']
        confname = request.data['confname']
        status = request.data['status']
        remark = request.data['remark']
        configcontent = request.data['options']
        customconfig = request.data['basicconfig']
        if id == '0':  # 新增数据
            chartbgcf = chartsbgconfig.objects.get_or_create(confname=confname,remark=remark,status=status,
                                   creater=request.user,configcontent = configcontent,customconfig=customconfig,
                                   createtime=datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                                   modifytime=datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
        else :
            try:
                sceneBackGroundConfigObj = chartsbgconfig.objects.get(id=id)
                sceneBackGroundConfigObj.confname = confname
                sceneBackGroundConfigObj.remark = remark
                sceneBackGroundConfigObj.status = status
                sceneBackGroundConfigObj.configcontent = configcontent
                sceneBackGroundConfigObj.customconfig = customconfig
                sceneBackGroundConfigObj.modifytime = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                sceneBackGroundConfigObj.save()
            except Exception as sceneConfGetError:
                resultObj = tools.errorMes('查询的信息不存在！')
    except Exception as errors:
        resultObj = tools.errorMes(errors.args)
    return Response(resultObj)

# 删除chartsbgconfig
@api_view(http_method_names=['GET'])
def deletechartsbgconfig(request):
    rs = {}
    id = request.GET['id']
    chartsbgconf = chartsbgconfig.objects.get(id=id)
    #获得项目当前路径，以供删除图片
    osdir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))).replace('\\', '/')
    if chartsbgconf:
        chartsbgconf.delete()
        #需要删除图片
        jsonConfig = chartsbgconf.configcontent
        if type(jsonConfig) == type('stringstringstring'):
            jsonConfig = json.loads(jsonConfig.replace("'", "\""))
        globalBGImage = ''
        if 'globalStyle' in jsonConfig and 'backgroundImage' in jsonConfig['globalStyle']:
            globalBGImage = osdir + os.sep + jsonConfig['globalStyle']['backgroundImage']
        chartBGImage = ''
        if 'chartStyle' in jsonConfig and 'backgroundImage' in jsonConfig['chartStyle']:
            chartBGImage = osdir + os.sep + jsonConfig['chartStyle']['backgroundImage']
        filterBGImage = ''
        if 'filterStyle' in jsonConfig and 'backgroundImage' in jsonConfig['filterStyle']:
            filterBGImage = osdir + os.sep + jsonConfig['filterStyle']['backgroundImage']
        alertMsg = ''
        if globalBGImage:
            try:
                os.remove(globalBGImage)
            except Exception as e:
                alertMsg += '删除全局配置的图片没找到！'
        if chartBGImage:
            try:
                os.remove(chartBGImage)
            except Exception as e:
                alertMsg += '删除组件背景配置的图片没找到！'
        if filterBGImage:
            try:
                os.remove(filterBGImage)
            except Exception as e:
                alertMsg += '删除过滤器北京配置的图片没找到！'
        rs['code'] = 0
        if alertMsg:
            rs['deleteMsg'] = alertMsg
    else:
        rs['code'] = 1
        rs['msg'] = '未找到删除数据！'
    return Response(rs)

# 是否禁用
@api_view(http_method_names=['GET'])
def forbiddenchartsbgconfig(request):
    resultObj = tools.successMes()
    try:
        id = request.GET.get("id")
        chartsbgObj = chartsbgconfig.objects.get(id=id)
        now_status = chartsbgObj.status
        if str(now_status) == '1':
            chartsbgObj.status = 0
        else:
            chartsbgObj.status = 1
        chartsbgObj.save()
    except Exception as e:
        resultObj = tools.errorMes(e.args)
    return Response(resultObj)

# 保存
@api_view(http_method_names=['POST'])
def savechartsbgconfig(request):  # 接收用户传入的数据，将css内容写入新的文件保存到制定目录下，将必要数据持久化到数据库
    rs = {}
    id = '0';  ##新增
    if 'id' in request.data:
        id = request.data['id']
    confname = request.data['confname']
    status = request.data['status']
    remark = request.data['remark']
    filecontent = request.data['filecontent']
    filename = str(int(time.time()) + random.randint(0, 9999)) + '.css'  # 产生filter_module_143473文件名，避免文件名重复
    # 以此文件名创建文件到指定文件夹下，把filecontent内容写入
    osdir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))).replace('\\', '/')
    themeChtCssDir = '/frontend/upload/user_default/'
    savefilepath = osdir + themeChtCssDir + filename
    # print('savefilepath=',savefilepath)

    with open(savefilepath, 'w') as cssfile:
        cssfile.write(filecontent)

    if id == '0':  # 新增数据

        chartbgcf = chartsbgconfig.objects.get_or_create(confname=confname, filename=filename,
                                                         fileurl=themeChtCssDir + filename
                                                         , filecontent=filecontent, remark=remark, status=status,
                                                         creater=request.user
                                                         , createtime=datetime.datetime.now().strftime(
                '%Y-%m-%d %H:%M:%S')
                                                         , modifytime=datetime.datetime.now().strftime(
                '%Y-%m-%d %H:%M:%S'))
        rs['code'] = 1
    else:  # 编辑数据
        chartbgcfg = chartsbgconfig.objects.get(id=id)
        if chartbgcfg:
            # 删除旧的css文件
            osdir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))).replace('\\', '/')
            furl = chartbgcfg.fileurl
            cssFilePath = osdir + furl
            if os.path.exists(cssFilePath):
                os.remove(cssFilePath)

            chartbgcfg.confname = confname
            chartbgcfg.filename = filename
            chartbgcfg.fileurl = themeChtCssDir + filename
            chartbgcfg.filecontent = filecontent
            chartbgcfg.remark = remark
            chartbgcfg.status = status
            chartbgcfg.modifytime = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            chartbgcfg.save()

            rs['code'] = 1
        else:
            rs['code'] = 0
            rs['msg'] = '没有此条数据!'

    return Response(rs)


# 校验
@api_view(http_method_names=['POST'])
def checkchartsbgconfname(request):
    confname = request.data['name']
    id = request.data['id']
    rs = {}
    try:
        cfname = chartsbgconfig.objects.filter(confname=confname)
        if not cfname:
            rs['code'] = 1
        else:
            if id != cfname[0].id:
                rs['code'] = 0
                rs['msg'] = '名称重复'
            else:
                rs['code'] = 1
    except Exception as e:
        rs['code'] = 0
        rs['msg'] = '异常错误'
    return Response(rs)

# 导出配置项
@api_view(http_method_names=['GET'])
def exportBGConfig(request):
    id = request.GET['id']
    rsObj = tools.successMes()
    try:
        cfname = chartsbgconfig.objects.get(id=id)
        copyObj = {} #拷贝的数据项，需要将此对象转为json
        copyObj['confname'] = cfname.confname
        copyObj['configcontent'] = cfname.configcontent
        copyObj['remark'] = cfname.remark
        copyObj['status'] = cfname.status
        copyObj['options'] = cfname.options
        copyObj['customconfig'] = cfname.customconfig
        copyObj['filecontent'] = cfname.filecontent
        copyObj['filename'] = cfname.filename
        copyObj['fileurl'] = cfname.fileurl
        #找到图片。场景的、组件的、过滤器的
        osdir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))).replace('\\', '/')#项目目录

        jsonConfig = cfname.configcontent
        if type(jsonConfig) == type('stringstringstring'):
            jsonConfig = json.loads(jsonConfig.replace("'","\""))
        globalBGImage = ''
        if 'globalStyle' in jsonConfig and 'backgroundImage' in jsonConfig['globalStyle']:
            globalBGImage = osdir + os.sep + jsonConfig['globalStyle']['backgroundImage']
        chartBGImage = ''
        if 'chartStyle' in jsonConfig and 'backgroundImage' in jsonConfig['chartStyle']:
            chartBGImage = osdir + os.sep + jsonConfig['chartStyle']['backgroundImage']
        filterBGImage = ''
        if 'filterStyle' in jsonConfig and 'backgroundImage' in jsonConfig['filterStyle']:
            filterBGImage = osdir + os.sep + jsonConfig['filterStyle']['backgroundImage']
        #在temp文件夹里创建一个样式配置的文件夹，将copyObj写入json文件中，将图片写入文件夹，压缩文件夹并提供js前端下载
        #如果第二次下载则需要删除第一次的文件夹和压缩文件
        #在temp文件夹的文件会被调度删除（每晚1点左右）
        timeStampSuff = str(int(time.time())) #增加时间戳,避免重复
        exportBGConfigParentDir = r'/frontend/upload/temp_files/sceneBGConfig_' + timeStampSuff + '/'
        exportBGConfigDir = osdir + exportBGConfigParentDir
        if os.path.exists(exportBGConfigDir):
            shutil.rmtree(exportBGConfigDir)#清空文件夹，但改文件夹也被清空了
        os.mkdir(exportBGConfigDir)#创建文件夹
        with open(exportBGConfigDir + 'config.json','w') as writer:#将json数据写入config.json
            json.dump(copyObj,writer)
        #将图片写入文件夹
        if globalBGImage :
            fileName = os.path.basename(globalBGImage)
            if fileName:#这里不判断上传文件的图片类型
                shutil.copy(globalBGImage,exportBGConfigDir + fileName)#将图片拷贝到exportBGConfigParentDir目录下
        if chartBGImage:
            fileName = os.path.basename(chartBGImage)
            if fileName:  # 这里不判断上传文件的图片类型
                shutil.copy(chartBGImage, exportBGConfigDir + fileName)  # 将图片拷贝到exportBGConfigParentDir目录下
        if filterBGImage:
            fileName = os.path.basename(filterBGImage)
            if fileName:  # 这里不判断上传文件的图片类型
                shutil.copy(filterBGImage, exportBGConfigDir + fileName)  # 将图片拷贝到exportBGConfigParentDir目录下
        #将文件压缩成zip文件并提供给js下载
        zipFileName = exportBGConfigDir[:-1] + '.zip'
        zf = zipfile.ZipFile(zipFileName,'w',zipfile.ZIP_DEFLATED)
        for dirPath,dirNames,fileNames in os.walk(exportBGConfigDir):
            # fPath = dirPath.replace(exportBGConfigDir,'')#需要替换根目录，不这样做就会从根目录复制
            fPath = 'sceneBGConfig_' + timeStampSuff + os.sep #需要替换根目录，不这样做就会从根目录复制
            for fileName in fileNames:
                zf.write(os.path.join(dirPath,fileName),fPath + fileName)
        logging.info('----场景配置数据包压缩成功-----')
        zf.close()

        rsObj['downLoadZipFilePath'] = exportBGConfigParentDir[:-1] + '.zip'    #将路径返回给js以供下载
        rsObj['timeStampSuff'] = timeStampSuff    #记录时间戳，前台js使用
    except Exception as e:
        rsObj = tools.errorMes('异常错误!:%s' % e.args)
    return Response(rsObj)


#递归遍历文件夹内的文件，返回根目录rootDir下所有文件名
def getAllFileNameByDir(rootDir):
    allFilesName = []
    for lists in os.listdir(rootDir):
        path = os.path.join(rootDir, lists)
        if os.path.isdir(path):
            allFilesName.extend(getAllFileNameByDir(path))
        else:
            allFilesName.append(path)
    return allFilesName

#根据文件名在文件路径下查找文件
def getFileName(filePath,fileName):
    if not filePath or not os.path.exists(filePath):
        raise Exception('文件路径不对！')
    if not fileName:
        raise Exception('文件名出错！')
    allFilesName = getAllFileNameByDir(filePath)#文件夹下所有文件名
    for fname in allFilesName:
        if fileName in fname:#遍历查找文件后缀名
            return fname
    return ''

# 导出配置项
@api_view(http_method_names=['POST'])
def uploadCoonfigFile(request):
    rsObj = tools.successMes()
    # 项目目录
    osdir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))).replace('\\', '/')
    parentPath = osdir + r'/frontend/upload/temp_files/' + str(int(time.time())) + '/'
    try:
        uploadedFile = request.FILES['filename']
        fileName = uploadedFile.name
        if not fileName.endswith('.zip'):
            raise Exception('上传文件类型出错，只能上传zip格式的文件！')

        if not os.path.exists(parentPath):
            os.makedirs(parentPath)
        savePath = parentPath + fileName
        if os.path.exists(savePath) and os.path.isfile(savePath):#如果目录已经存在就删除
            shutil.rmtree(savePath)
        with open(savePath, 'wb+') as destination:
            for chunk in uploadedFile.chunks():
                destination.write(chunk)
        #到此为止文件上传成功，存入到/frontend/upload/temp_files/下了
        #开始解压文件
        if not zipfile.is_zipfile(savePath):
            raise Exception('上传文件类型出错，只能上传zip文件！')
        with zipfile.ZipFile(savePath) as f:
            for fn in f.namelist():
                f.extract(fn,parentPath)
            f.close()
        #解析解压json文件
        configJsonFileName = getFileName(parentPath,'.json')
        #解析json文件，确认json里所需的图片是否都存在，如果都存在就保存到本地，否则抛出异常
        if not configJsonFileName:
            raise Exception('查找JSON文件出错，请检查上传文件！')
        # print('----configJsonFileName=',configJsonFileName)

        configJsonAllObj = {}#解析json对象
        with open(configJsonFileName,'r') as reader:
            configJsonAllObj = json.load(reader)
        if not configJsonAllObj:
            raise Exception('JSON解析出错！')

        pictureList = []#找到所有的图片并把图片保存到本地
        configJsonObj = configJsonAllObj['configcontent'] if 'configcontent' in configJsonAllObj else {}
        if type(configJsonObj) == type('stringstringstring'):
            configJsonObj = json.loads(configJsonObj )
        if 'globalStyle' in configJsonObj:
            if 'backgroundImage' in configJsonObj['globalStyle']:
                pictureList.append(osdir + configJsonObj['globalStyle']['backgroundImage'])
        if 'chartStyle' in configJsonObj:
            if 'backgroundImage' in configJsonObj['chartStyle']:
                pictureList.append(osdir + configJsonObj['chartStyle']['backgroundImage'])
        if 'filterStyle' in configJsonObj:
            if 'backgroundImage' in configJsonObj['filterStyle']:
                pictureList.append(osdir + configJsonObj['filterStyle']['backgroundImage'])

        #必要字段检查
        if not configJsonAllObj['confname'] or not configJsonAllObj['configcontent']:
            raise Exception('上传JSON数据有误！请检查后重新上传！')
        # 将图片文件拷贝到指定文件夹，如果文件不存在就抛出异常
        uploadPicFiles = []#所有能拷贝的图片（从上传中解析出来），如果图片寻找中错了则需要抛出异常，即便错了一张也需要抛出异常
        for picFile in pictureList:
            picFileJsonName = os.path.basename(picFile)#获取json数据的图片名称
            uploadPicFile = getFileName(parentPath,picFileJsonName)#从上传文件中查找是否有该json里的picture的名字，如果有就保存
            if not uploadPicFile:
                raise Exception('上传图片名和JSON中的图片名不一致或图片没有！')
            uploadPicFiles.append(uploadPicFile)
        for picFile in pictureList:
            for uploadPic in uploadPicFiles:
                if os.path.basename(picFile) == os.path.basename(uploadPic):#把图片从上传解压的文件夹拷贝到json里对应的文件夹
                    if not os.path.exists(os.path.dirname(picFile)):
                        os.makedirs(os.path.dirname(picFile))
                    shutil.copy(uploadPic,picFile)#拷贝

        #将对象持久化到数据库
        uploadUserName = request.user
        chartsbgconfig.objects.create(confname=configJsonAllObj['confname'],
                                      configcontent=configJsonAllObj['configcontent'],
                                      remark=configJsonAllObj['remark'],
                                      creater=uploadUserName,
                                      status=configJsonAllObj['status'],
                                      options=configJsonAllObj['options'],
                                      customconfig=configJsonAllObj['customconfig'],
                                      filecontent=configJsonAllObj['filecontent'],
                                      filename=configJsonAllObj['filename'],
                                      fileurl=configJsonAllObj['fileurl'])


        rsObj['data'] = '上传成功'
    except Exception as e:
        rsObj = tools.errorMes('异常错误!:%s' % e)
    finally:
        # 删除上传到temp的文件
        shutil.rmtree(parentPath)

    return Response(rsObj)

# 传入scene的htmlcleanconfig解析，返回chatsid[]
def resolveHtml(htmlstr):
    # chartsids={}
    chartslist = []
    # tableslist=[]
    bsObj = BeautifulSoup(htmlstr, 'html.parser')
    chartscontent = bsObj.find_all("div", {'class': 'chart-content'})
    for chart in chartscontent:
        dtype = chart.attrs['data-type']
        if dtype == 'echarts':
            chartid = chart.find_all_next("span", {"class": "ecdatatag"})[0].attrs['id']
            if chartid.index('_') > 0:
                chartid = chartid[chartid.index('_') + 1:]
                chartslist.append(chartid)
        elif dtype == 'table':
            chartid = chart.find_all_next("span", {"class": "tdatatag"})[0].attrs['id']
            if chartid.index('_') > 0:
                chartid = chartid[chartid.index('_') + 1:]
                chartid = chartid + "_t"  # 区别于chart的标识
                chartslist.append(chartid)

    # chartsids['chartslist']=chartslist
    # chartsids['tableslist']=tableslist
    # print('chartsids=',chartsids)
    return chartslist


# 获取所有导出主题、场景、组件的id数据
@api_view(http_method_names=['GET'])
def getExportAllData(request):
    expalldata = {}
    themelist = []
    themesobj = themes.objects.all()  # 需要添加过滤条件再根据情况添加
    for theme in themesobj:
        obj = {}  # 可扩展传递给页面的数据
        obj['check'] = False
        obj['id'] = theme.id
        obj['name'] = theme.name
        obj['scenesconfig'] = theme.scenesconfig
        themelist.append(obj)
    scenelist = []
    scenesobj = scenes.objects.all()  # 需要添加过滤条件再根据情况添加
    for scene in scenesobj:
        obj = {}  # 可扩展传递给页面的数据
        obj['check'] = False
        obj['id'] = scene.id
        obj['name'] = scene.name
        obj['chartsid'] = resolveHtml(scene.htmlcleanconfig)
        scenelist.append(obj)
    chartlist = []
    chartsobj = charts.objects.all()  # 需要添加过滤条件再根据情况添加
    for chart in chartsobj:
        obj = {}  # 可扩展传递给页面的数据
        obj['check'] = False
        obj['id'] = chart.id
        obj['name'] = chart.name
        obj['datasetstring'] = chart.datasetstring
        chartlist.append(obj)
    datatablesobj = datatable.objects.all()  # 需要添加过滤条件再根据情况添加
    for tb in datatablesobj:
        obj = {}  # 可扩展传递给页面的数据
        obj['check'] = False
        obj['id'] = tb.id + "_t"
        obj['name'] = tb.name
        obj['datasetstring'] = tb.jsonconfig  # table的olap数据存储再jsconfig里，为了和chart统一，存放在datasetstring里
        chartlist.append(obj)

    expalldata['themelist'] = themelist
    expalldata['scenelist'] = scenelist
    expalldata['chartlist'] = chartlist
    # print('json.dump(expalldata)=',json.dumps(expalldata))
    return Response(json.dumps(expalldata))


@api_view(http_method_names=['GET'])
def searchExpByType(request):  # 通过type来判断查找类型，通过searchkey来模糊查询
    type = request.GET['type']
    searchkey = request.GET['search']
    if type == 'theme':
        themelist = []
        if searchkey:
            themesobj = themes.objects.filter(name__contains=searchkey)  # 需要添加过滤条件再根据情况添加
        else:
            themesobj = themes.objects.all()  # 需要添加过滤条件再根据情况添加
        for theme in themesobj:
            obj = {}  # 可扩展传递给页面的数据
            obj['check'] = False
            obj['id'] = theme.id
            obj['name'] = theme.name
            obj['scenesconfig'] = theme.scenesconfig
            themelist.append(obj)
        return Response(themelist)
    if type == 'scene':
        scenelist = []
        if searchkey:
            scenesobj = scenes.objects.filter(name__contains=searchkey)  # 需要添加过滤条件再根据情况添加
        else:
            scenesobj = scenes.objects.all()  # 需要添加过滤条件再根据情况添加
        for scene in scenesobj:
            obj = {}  # 可扩展传递给页面的数据
            obj['check'] = False
            obj['id'] = scene.id
            obj['name'] = scene.name
            obj['chartsid'] = resolveHtml(scene.htmlcleanconfig)
            scenelist.append(obj)
        return Response(scenelist)
    if type == 'chart':
        chartlist = []
        if searchkey:
            chartsobj = charts.objects.filter(name__contains=searchkey)  # 需要添加过滤条件再根据情况添加
        else:
            chartsobj = charts.objects.all()  # 需要添加过滤条件再根据情况添加
        for chart in chartsobj:
            obj = {}  # 可扩展传递给页面的数据
            obj['check'] = False
            obj['id'] = chart.id
            obj['name'] = chart.name
            obj['datasetstring'] = chart.datasetstring
            chartlist.append(obj)
        if searchkey:
            datatablesobj = datatable.objects.filter(name__contains=searchkey)  # 需要添加过滤条件再根据情况添加
        else:
            datatablesobj = datatable.objects.all()  # 需要添加过滤条件再根据情况添加
        for tb in datatablesobj:
            obj = {}  # 可扩展传递给页面的数据
            obj['check'] = False
            obj['id'] = tb.id + "_t"
            obj['name'] = tb.name
            obj['datasetstring'] = tb.jsonconfig  # table的olap数据存储再jsconfig里，为了和chart统一，存放在datasetstring里
            chartlist.append(obj)
        return Response(chartlist)
    return Response([])  # 如果没匹配上就返回空


# def model_to_dict_byMyself(modelobj,id):#重写model_to_dict方法，加上id字段
#     modeldict=model_to_dict(modelobj)
#     print('modeldict=',modeldict)
#     print('modelobj=',modelobj)
#     modeldict['id']=id
#     return modeldict

def expTableObj(themelist, scenelist, chartlist):
    sceneidlist = [sceneo['id'] for sceneo in scenelist]
    chartidlist = [charto['id'] for charto in chartlist]
    for themelst in themelist:
        for sceneFrTheme in json.loads(themelst['scenesconfig']):  # 把scenelist中不存在themelist子节点的重新加入到scenelist
            if sceneFrTheme['id'] not in sceneidlist:
                sceneidlist.append(sceneFrTheme['id'])
    for scenelst in scenelist:
        for chartFrScene in scenelst['chartsid']:  # 同上
            if chartFrScene not in chartidlist:
                chartidlist.append(chartFrScene)
    # 如果只选择了主题，没有选择场景和chart，需要根据主题找到场景，再根据场景找到chart
    if len(themelist) and not len(scenelist):
        for scenenid in sceneidlist:
            scenemd = scenes.objects.get(id=scenenid)
            chartidlist.extend(resolveHtml(scenemd.htmlcleanconfig))

    # print('sceneidlist=',sceneidlist)
    # print('chartidlist=',chartidlist)
    themeTbObj = []
    sceneTbObj = []
    chartTbObj = []
    datatableTbObj = []
    olapTbObj = []
    chartTypeIds = []  # 先将所有的charttypeid全部导入，然后转成set就去除了重复值
    for temptheme in themelist:
        themeTbObj.append(utils.model_to_dict_wrapper(themes.objects.get(id=temptheme['id'])))
    for tempsceneid in sceneidlist:
        sceneTbObj.append(utils.model_to_dict_wrapper(scenes.objects.get(id=tempsceneid)))
        chartTypeIds.append(scenes.objects.get(id=tempsceneid).kind)
    distinguishOlapId = []
    for tempchartid in chartidlist:
        if tempchartid.endswith('_t'):
            tempchartid = tempchartid[0:tempchartid.index('_')]
            dtobj = datatable.objects.get(id=tempchartid)
            datatableTbObj.append(utils.model_to_dict_wrapper(dtobj))
            chartTypeIds.append(dtobj.kind)  # charttype
            dtobjconfig = eval(dtobj.jsonconfig)
            if dtobjconfig['olapid'] not in distinguishOlapId:  # 去重，chart和table可能公用一个olap
                distinguishOlapId.append(dtobjconfig['olapid'])
                olapTbObj.append(utils.model_to_dict_wrapper(olap.objects.get(id=dtobjconfig['olapid'])))
                chartTypeIds.append(olap.objects.get(id=dtobjconfig['olapid']).charttype)  # charttype
        else:
            chobj = charts.objects.get(id=tempchartid)
            chartTbObj.append(utils.model_to_dict_wrapper(chobj))
            chartTypeIds.append(chobj.kind)  # charttype
            chobjconf = eval(chobj.datasetstring)
            if chobjconf['id'] not in distinguishOlapId:  # 去重，chart和table可能公用一个olap
                distinguishOlapId.append(chobjconf['id'])
                olapTbObj.append(utils.model_to_dict_wrapper(olap.objects.get(id=chobjconf['id'])))
                chartTypeIds.append(olap.objects.get(id=chobjconf['id']).charttype)  # charttype
    # get charttypeobj
    chartTypeObj = []
    for charttypeo in list(set(chartTypeIds)):
        chartTypeObj.append(utils.model_to_dict_wrapper(charttype.objects.get(id=charttypeo)))

    rsobj = {}
    rsobj['themeTbObj'] = themeTbObj
    rsobj['sceneTbObj'] = sceneTbObj
    rsobj['chartTbObj'] = chartTbObj
    rsobj['datatableTbObj'] = datatableTbObj
    rsobj['olapTbObj'] = olapTbObj
    rsobj['chartTypeObj'] = chartTypeObj
    return rsobj


def getTbObjFrOlapID(olaplist):
    olapcolumnls = []
    olapextcolsls = []
    for op in olaplist:
        for olapcol in olapcolumn.objects.filter(olapid=op['id']):
            olapcolumnls.append(utils.model_to_dict_wrapper(olapcol))
        for olapext in olapextcols.objects.filter(olapid=op['id']):
            olapextcolsls.append(utils.model_to_dict_wrapper(olapext))
    return olapcolumnls, olapextcolsls


def getOlapTableAndData(olaplist, olapnum):  # 根据olap，导出对应的表和表的数据，表的数据顺序按照olapcolumn里存的顺序
    alltabledata = []
    for op in olaplist:
        tbobj = {}
        # 表名
        tbobj['tablename'] = op['name']
        # 导出表结构
        # tbstruct=[]
        # tbstructdt = utils.getCreateTableStructure(op['name'], utils.DATAXEXTENSION_DB_CHAR)
        tbstructdt = utils.SqlUtils(dbType=DBTypeCode.EXTENSION_DB.value).getCreateTableStructure(op['name'])
        # for dt in tbstructdt:
        #     print('dt=',dt)
        #     rows={}
        #     rows['table_schema']=dt[0]
        #     rows['table_name']=dt[1]
        #     rows['column_name']=dt[2]
        #     rows['data_type']=dt[3]
        #     rows['column_default']=dt[4]
        #     rows['is_nullable']=dt[5]
        #     tbstruct.append(rows)
        tbobj['tablestruct'] = tbstructdt
        # 导出表数据
        tbcols = []
        for dt in tbstructdt:  # 获取表字段信息
            tbcols.append(dt['columnName'])
        # 根据tbcols的顺序(列顺序)查询值
        # print('get tabledata sql=','select '+','.join(tbcols)+' from '+ op['name'] +' limit '+str(olapnum))
        tableData = []
        # tabledt = utils.getResultBySql('select ' + ','.join(tbcols) + ' from ' + op['name'] + ' limit ' + str(olapnum), utils.DATAXEXTENSION_DB_CHAR)
        tabledt = utils.SqlUtils(DBTypeCode.EXTENSION_DB.value).getArrResultWrapper('select ' + ','.join(tbcols) + ' from ' + op['name'] + ' limit ' + str(olapnum),logger,'dashboardviews.py', 'getOlapTableAndData')
        realcount = 0
        # 模糊数据
        for tbdata in tabledt:
            # tableData.append(dict(zip(tbcols,tbdata)))#把查询出来的数据根据列名合并成字典
            factor = random.random()
            sigleRow = {}
            i = 0
            for tbcol in tbcols:
                sigleRow[tbcol] = tbdata[i] * factor if isinstance(tbdata[i], float) else tbdata[i]  # 只对float数据进行脱敏处理
                i += 1
            tableData.append(sigleRow)
            realcount += 1
        tbobj['realcount'] = realcount  # 记录导出数据实际条数
        tbobj['tabledata'] = tableData

        alltabledata.append(tbobj)

    return alltabledata


class CJsonEncoder(json.JSONEncoder):  # 扩展json.dumps能够转换时间类型
    def default(self, obj):
        if isinstance(obj, datetime.datetime):
            return obj.strftime('%Y-%m-%d %H:%M:%S')
        elif isinstance(obj, datetime.date):
            return obj.strftime("%Y-%m-%d %H:%M:%S")
        elif isinstance(obj, date):
            return obj.strftime("%Y-%m-%d")
        else:
            return json.JSONEncoder.default(self, obj)


@api_view(http_method_names=['POST'])
def exportPackage(request):
    rs = {}
    try:
        themelist = request.data['themels']
        scenelist = request.data['scenels']
        chartlist = request.data['chartls']  # 有table和chart的id，table的id后面有_t标识
        olapnum = request.data['olapnum']
        # 根据传入的参数，查询list，导出结果为json格式obj={}，里面包含了theme，scene、chart、connect_olap、connect_column、connect_olapextcols
        rsobj = {}
        packageobj = expTableObj(themelist, scenelist, chartlist)  # 根据id查询数据库表的数据存放到rsobj里
        # print('packageobj=',packageobj)
        rsobj['themels'] = packageobj['themeTbObj']
        rsobj['scenels'] = packageobj['sceneTbObj']
        rsobj['chartls'] = packageobj['chartTbObj']
        rsobj['datatablels'] = packageobj['datatableTbObj']
        rsobj['olapls'] = packageobj['olapTbObj']
        rsobj['charttypels'] = packageobj['chartTypeObj']
        # print('rsobj0', json.dumps(rsobj))
        olapcolumnsls, olapextcolsls = getTbObjFrOlapID(rsobj['olapls'])  # 通过olapid获取olapcolumn和olapextcols
        rsobj['olapcolumnsls'] = olapcolumnsls
        rsobj['olapextcolsls'] = olapextcolsls
        # print('rsobj1',json.dumps(rsobj))
        # 根据olap，导出对应的表结构和表的数据，表的数据顺序按照olapcolumn里存的顺序
        rsobj['alltableanddata'] = getOlapTableAndData(rsobj['olapls'], olapnum)
        # print('rsobj2', json.dumps(rsobj,cls=CJsonEncoder))

        # 将rsobj转成json并保存到目录下，返回url给页面让用户下载
        saveFileName = datetime.datetime.now().strftime('%Y-%m-%d_%H_%M_%S') + '.json'
        saveDir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        # suffDir='/frontend/upload/temp_files'
        suffDir = os.path.join('frontend', 'upload', 'temp_files')
        # savePath=saveDir+suffDir+os.sep+saveFileName
        tempfiledir = os.path.join(saveDir, 'frontend', 'upload', 'temp_files')
        savePath = tempfiledir + os.sep + saveFileName
        savePath = savePath.replace('\\', '/')
        # print('savePath=',savePath)
        with open(savePath, 'w') as f:
            f.write(json.dumps(rsobj, cls=CJsonEncoder))
        rs['code'] = 1
        rs['filepath'] = (os.sep + suffDir + os.sep + saveFileName).replace('\\', '/')
        rs['filename'] = saveFileName
        # print('rs=',rs)
        ############把导出信息输出到“控制台（页面上的textarea）”
        output2console = {}
        # theme导出记录
        exportThemeNames = [themeo['name'] for themeo in themelist]
        output2console['exportThemeNames'] = exportThemeNames
        # scene导出记录，从expTableObj的结果中获取
        exportSceneNames = []
        for tmpobj in rsobj['scenels']:
            exportSceneNames.append(tmpobj['name'])
        output2console['exportSceneNames'] = exportSceneNames
        # chart导出记录
        exportChartNames = []
        for tmpobj in rsobj['chartls']:
            exportChartNames.append(tmpobj['name'])
        output2console['exportChartNames'] = exportChartNames
        exportTableChartNames = []
        for tmpobj in rsobj['datatablels']:
            exportTableChartNames.append(tmpobj['name'])
        output2console['exportTableChartNames'] = exportTableChartNames
        # table(olap)导出记录
        exportTablesInfo = []
        for tmpobj in rsobj['alltableanddata']:
            tempDict = {}
            tempDict['tablename'] = tmpobj['tablename']
            tempDict['tabledatarealnum'] = tmpobj['realcount']
            exportTablesInfo.append(tempDict)
        output2console['exportTablesInfo'] = exportTablesInfo
        rs['outputconsole'] = json.dumps(output2console)
    except Exception as e:
        raise e
        rs['code'] = 0
        rs['msg'] = e.args
    return Response(rs)


# 导入配置包,查询所有数据
@api_view(http_method_names=['GET'])
def getImportAllData(request):
    rs = {}
    try:
        impConfModels = []
        if 'searchkey' in request.GET:
            impConfModels = impconfpackage.objects.filter(filename=request.GET['searchkey']).order_by('-createdate')
        else:
            impConfModels = impconfpackage.objects.all().order_by('-createdate')  # 降序排列

        allImpConfData = []
        for obj in impConfModels:
            dt = {}
            dt['id'] = obj.id
            dt['fileName'] = obj.filename
            dt['fileUrl'] = obj.fileurl
            dt['creater'] = obj.creater
            dt['createDate'] = (obj.createdate).strftime("%Y-%m-%d %H:%M:%S")
            allImpConfData.append(dt)
        rs['code'] = 1
        rs['data'] = allImpConfData
    except Exception as e:
        rs['code'] = 0
        rs['msg'] = e.args
    return Response(rs)


def getAndInsertDtFrFile(fileUrl):
    """解析fileUrl下的文件，将解析数据插入数据库，搜集插入数据信息并返回"""
    # print('getAndInsertDtFrFile',fileUrl)
    insertLogs = {}  # 日志，返回到页面显示给用户
    insertLogs['themelogs'] = []
    insertLogs['scenelogs'] = []
    insertLogs['chartlogs'] = []
    insertLogs['datatablelogs'] = []
    insertLogs['olaplogs'] = []
    insertLogs['olapcolumnlogs'] = []
    insertLogs['olapextlogs'] = []
    insertLogs['olaptablelogs'] = []
    theImpIds = {}  # 将插入的数据的id记录下来，当要删除时需要使用
    theImpIds['themeids'] = []
    theImpIds['sceneids'] = []
    theImpIds['chartids'] = []
    theImpIds['datatableids'] = []
    theImpIds['olapids'] = []
    theImpIds['olapcolumnids'] = []
    theImpIds['olapextids'] = []
    theImpIds['charttypeids'] = []
    theImpIds['olaptable'] = []
    if os.path.exists(fileUrl):
        with open(fileUrl, encoding='utf-8') as fl:
            allData = json.load(fl)
        # print('allData=',allData)
        # 将数据插入theme，scene，chart，datatable，olap，olapcolumn，olapext，新建表，插入表数据
        impThemeData = []
        if 'themels' in allData:
            impThemeData = allData['themels']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = '主题为空！'
        for themeo in impThemeData:
            themes.objects.get_or_create(id=themeo['id'], name=themeo['name']
                                         , code=themeo['code'], kind=themeo['kind']
                                         , scenesconfig=themeo['scenesconfig'], themeconfig=themeo['themeconfig']
                                         , remark=themeo['remark'], createname=themeo['createname']
                                         , createtime=themeo['createtime'], orderby=themeo['orderby']
                                         , refreshspeed=themeo['refreshspeed'], continues=themeo['continues']
                                         , keywords=themeo['keywords'], allowusers=themeo['allowusers']
                                         , options=themeo['options']
                                         )
            theImpIds['themeids'].append(themeo['id'])
            insertLogs['themelogs'].append(themeo['name'])
        # 将数据插入到scene
        impSceneData = []
        if 'scenels' in allData:
            impSceneData = allData['scenels']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = '场景为空！'
        for sceneo in impSceneData:
            scenes.objects.get_or_create(id=sceneo['id'], name=sceneo['name']
                                         , kind=sceneo['kind'], htmlrawstring=sceneo['htmlrawstring']
                                         , htmlcleanconfig=sceneo['htmlcleanconfig'], remark=sceneo['remark']
                                         , createname=sceneo['createname'], createtime=sceneo['createtime']
                                         , orderby=sceneo['orderby'], refreshspeed=sceneo['refreshspeed']
                                         , keywords=sceneo['keywords'], allowusers=sceneo['allowusers']
                                         , basicconfig=sceneo['basicconfig'], options=sceneo['options']
                                         )
            theImpIds['sceneids'].append(sceneo['id'])
            insertLogs['scenelogs'].append(sceneo['name'])
        # 将数据插入chart表
        impChartsData = []
        if 'chartls' in allData:
            impChartsData = allData['chartls']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = '图表组件为空！'
        for charto in impChartsData:
            charts.objects.get_or_create(id=charto['id'], name=charto['name']
                                         , kind=charto['kind'], jsonconfig=charto['jsonconfig']
                                         , remark=charto['remark'], createname=charto['createname']
                                         , createtime=charto['createtime'], orderby=charto['orderby']
                                         , filterstring=charto['filterstring'], datasetstring=charto['datasetstring']
                                         , charttype=charto['charttype'], echartconfig=charto['echartconfig']
                                         , refreshspeed=charto['refreshspeed'], keywords=charto['keywords']
                                         , allowusers=charto['allowusers'], imgpath=charto['imgpath']
                                         , options=charto['options']
                                         )
            theImpIds['chartids'].append(charto['id'])
            insertLogs['chartlogs'].append(charto['name'])
        # 将数据插入datatable表
        impDataTableData = []
        if 'datatablels' in allData:
            impDataTableData = allData['datatablels']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = '表格组件为空！'
        for datatableo in impDataTableData:
            options = {}
            if datatableo['options']:
                options = json.loads(datatableo['options'])
            options['isImportData'] = True
            datatable.objects.get_or_create(id=datatableo['id'], name=datatableo['name']
                                            , kind=datatableo['kind'], olapid=datatableo['olapid']
                                            , jsonconfig=datatableo['jsonconfig'], remark=datatableo['remark']
                                            , refreshspeed=datatableo['refreshspeed'], keywords=datatableo['keywords']
                                            , allowusers=datatableo['allowusers'], imgpath=datatableo['imgpath']
                                            , createname=datatableo['createname'], createtime=datatableo['createtime']
                                            , options=json.dumps(options)
                                            )
            theImpIds['datatableids'].append(datatableo['id'])
            insertLogs['datatablelogs'].append(datatableo['name'])
        # 将数据插入olap表
        impOlapData = []
        if 'olapls' in allData:
            impOlapData = allData['olapls']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = 'olap为空！'
        for olapo in impOlapData:
            olap.objects.get_or_create(id=olapo['id'], sourceid=olapo['sourceid']
                                       , name=olapo['name'], charttype=olapo['charttype']
                                       , desc=olapo['desc'], columns=olapo['columns']
                                       , filters=olapo['filters'], table=olapo['table']
                                       , status=olapo['status'], dispatchid=olapo['dispatchid']
                                       , businesstype=olapo['businesstype'], expand=olapo['expand']
                                       , ifexpand=olapo['ifexpand'], olaptype=olapo['olaptype']
                                       , dispatchconfig=olapo['dispatchconfig'], create_date=olapo['create_date']
                                       , modify_date=olapo['modify_date'], enabled=olapo['enabled']
                                       , tag_config=olapo['tag_config'], options=olapo['options']
                                       )
            theImpIds['olapids'].append(olapo['id'])
            insertLogs['olaplogs'].append(olapo['name'])
        # 将数据插入olapcolumns表
        impOlapColumnData = []
        if 'olapcolumnsls' in allData:
            impOlapColumnData = allData['olapcolumnsls']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = 'olapcolumnsls为空！'
        for olapcol in impOlapColumnData:
            olapcolumn.objects.get_or_create(id=olapcol['id'], olapid=olapcol['olapid']
                                             , column=olapcol['column'], title=olapcol['title']
                                             , options=olapcol['options']
                                             )
            theImpIds['olapcolumnids'].append(olapcol['id'])
            # insertLogs['olapcolumnlogs'].append(olapcol['name'])
        # 将数据插入olapext表
        impOlapExtColData = []
        if 'olapextcolsls' in allData:
            impOlapExtColData = allData['olapextcolsls']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = 'olapextcolsls为空！'
        for olapexto in impOlapExtColData:
            olapextcols.objects.get_or_create(id=olapexto['id'], olapid=olapexto['olapid']
                                              , title=olapexto['title'], name=olapexto['name']
                                              , coltype=olapexto['coltype'], configs=olapexto['configs']
                                              , create_date=olapexto['create_date'], modify_date=olapexto['modify_date']
                                              , options=olapexto['options']
                                              )
            theImpIds['olapextids'].append(olapexto['id'])
            # insertLogs['olapextlogs'].append(olapexto['name'])
        # 将数据插入charttype表
        impCharTypeData = []
        if 'charttypels' in allData:
            impCharTypeData = allData['charttypels']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = 'charttypels为空！'
        for charttypeo in impCharTypeData:
            charttype.objects.get_or_create(id=charttypeo['id'], type_name=charttypeo['type_name']
                                            , level=charttypeo['level'], parent_id=charttypeo['parent_id']
                                            , orderby=charttypeo['orderby'], status=charttypeo['status']
                                            , create_time=charttypeo['create_time'],
                                            motify_time=charttypeo['motify_time']
                                            , create_user=charttypeo['create_user'], options=charttypeo['options'])
            theImpIds['charttypeids'].append(charttypeo['id'])
        # 新建olap表并插入数据
        impTableStructAndData = []
        if 'alltableanddata' in allData:
            impTableStructAndData = allData['alltableanddata']
        else:
            insertLogs['code'] = 0
            insertLogs['msg'] = '表格和数据为空！'

        for tdata in impTableStructAndData:
            createdTbLogs = {}
            tname = tdata['tablename']
            tstruct = tdata['tablestruct']
            tdatacount = tdata['realcount']
            tdt = tdata['tabledata']
            # createInfoLogs = utils.createTableFromStruct(tname, tstruct,utils.DATAXEXTENSION_DB_CHAR)  # 创建表结构
            createInfoLogs = utils.SqlUtils(DBTypeCode.EXTENSION_DB.value).createTableFromStruct(tname, tstruct)  # 创建表结构

            theImpIds['olaptable'].append(tname)
            if 'code' in createInfoLogs and createInfoLogs['code'] == 1:
                createdTbLogs['tablename'] = tname
                # 执行数据插入
                if 'tbFields' in createInfoLogs:
                    insertIntoTbLogs = utils.SqlUtils(DBTypeCode.EXTENSION_DB.value).insertIntoTable(tname, createInfoLogs['tbFields'], tdt,logger)  # 插入表数据
                    insertLogs['code'] = insertIntoTbLogs['code']
                    insertLogs['msg'] = insertIntoTbLogs['msg']
                    createdTbLogs['realDtCount'] = tdatacount
                else:
                    createdTbLogs['realDtCount'] = 0
                    insertLogs['code'] = 0
                    insertLogs['msg'] = '基础字段为空！'
            else:
                createdTbLogs['tablename'] = createInfoLogs['msg'] if 'msg' in createInfoLogs else ''
                createdTbLogs['realDtCount'] = 0
                insertLogs['code'] = 0
                insertLogs['msg'] = 'olap表的表结构创建失败！'
            insertLogs['olaptablelogs'].append(createdTbLogs)

        # 记录所有日志并返回到页面显示
    else:
        insertLogs['code'] = 0
        insertLogs['errorMsg'] = '文件不存在！'

    return insertLogs, theImpIds


@api_view(http_method_names=['POST'])
def uploadPkgConf(request):
    rs = {}
    f = request.FILES['filename']
    uploadfilename = f.name
    createrNm = request.user.username
    fileSuff = uploadfilename[uploadfilename.index('.', len(
        uploadfilename) - 6) + 1:]  # 解析文件格式，json文件类型。如果uploadfilename=a.b.txt则解析不到文件类型
    createdate = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    # 保存文件到/frontend/upload/temp_files
    saveDir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    writeFileName = datetime.datetime.now().strftime("%Y-%m-%d_%H_%M_%S")  # 写入到本地的文件名
    fileUrl = os.path.join('frontend', 'upload', 'packageConfig') + os.sep + writeFileName + '.' + fileSuff  # 保存到数据库
    pgConfFiledir = os.path.join(saveDir, 'frontend', 'upload', 'packageConfig')
    savePath = pgConfFiledir + os.sep + writeFileName + '.' + fileSuff
    print('savePath=', savePath)
    with open(savePath, 'wb') as fw:  # 写入到指定目录
        fw.write(f.read())
    try:
        # 解析文件，插入数据
        rs, impIds = getAndInsertDtFrFile(savePath)
        # 保存到数据库
        impconfpackage.objects.get_or_create(filename=uploadfilename, fileurl=fileUrl, creater=createrNm,
                                             createdate=createdate, impIdsAndTbName=json.dumps(impIds))
    except Exception as e:
        rs['code'] = 0
        rs['msg'] = e.args
        raise e
    return Response(rs)


@api_view(http_method_names=['GET'])
def delUploadConf(request):
    print('request.GET=', request.GET)
    id = request.GET['id']
    print('id=', id)
    rs = {}
    # 删除表数据，删除表，删除文件
    impConfPkObj = impconfpackage.objects.get(id=id)
    if impConfPkObj.impIdsAndTbName:
        impIds = json.loads(impConfPkObj.impIdsAndTbName)
    else:
        impIds = {}
    print('impIds=', impIds)
    if 'themeids' in impIds and impIds['themeids']:  # 选择themeIds
        themeids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['themeids']))
        themes.objects.extra(where=['id in (' + themeids + ')']).delete()  # 过滤并删除
    if 'sceneids' in impIds and impIds['sceneids']:  # 选择Ids
        sceneids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['sceneids']))
        scenes.objects.extra(where=['id in (' + sceneids + ')']).delete()  # 过滤并删除
    if 'chartids' in impIds and impIds['chartids']:  # 选择Ids
        chartids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['chartids']))
        charts.objects.extra(where=['id in (' + chartids + ')']).delete()  # 过滤并删除
    if 'datatableids' in impIds and impIds['datatableids']:  # 选择Ids
        datatableids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['datatableids']))
        datatable.objects.extra(where=['id in (' + datatableids + ')']).delete()  # 过滤并删除
    if 'olapids' in impIds and impIds['olapids']:  # 选择Ids
        olapids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['olapids']))
        olap.objects.extra(where=['id in (' + olapids + ')']).delete()  # 过滤并删除
    if 'olapcolumnids' in impIds and impIds['olapcolumnids']:  # 选择Ids
        olapcolumnids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['olapcolumnids']))
        olapcolumn.objects.extra(where=['id in (' + olapcolumnids + ')']).delete()  # 过滤并删除
    if 'olapextids' in impIds and impIds['olapextids']:  # 选择Ids
        olapextids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['olapextids']))
        olapextcols.objects.extra(where=['id in (' + olapextids + ')']).delete()  # 过滤并删除
    if 'charttypeids' in impIds and impIds['charttypeids']:  # 选择Ids
        charttypeids = ','.join(map(lambda x: "\'" + str(x) + "\'", impIds['charttypeids']))
        charttype.objects.extra(where=['id in (' + charttypeids + ')']).delete()  # 过滤并删除

    session = utils.SqlUtils(DBTypeCode.EXTENSION_DB.value)
    try:
        # 删除表olap表
        if 'olaptable' in impIds and impIds['olaptable']:  # 选择Ids
            for olaptb in impIds['olaptable']:
                delOlapTableSql = 'DROP TABLE IF EXISTS "public"."' + olaptb + '";'
                session.executeUpdateSql(delOlapTableSql)
        session.closeConnect()
    except Exception as error:
        session.rollBack()
        logger.error('---error---file:dashboardviews.py;method:delUploadConf;error=%s' % error)
    # 删除文件
    saveDir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    pgConfFiledir = os.path.join(saveDir, impConfPkObj.fileurl)
    print('delete pgConfFiledir=', pgConfFiledir)
    if pgConfFiledir.endswith('.json') and os.path.exists(pgConfFiledir):
        os.remove(pgConfFiledir)
        rs['code'] = 1
    else:
        rs['code'] = 0
        rs['msg'] = '文件不存在！'
    impConfPkObj.delete()  # 删除列表记录

    return Response(rs)


# 获取所有数据标签
@api_view(http_method_names=['GET'])
def getDataPerTag(request):
    where = ''
    if 'search' in request.GET:
        where = where + """  where name like '%""" + request.GET['search'] + """%' """
    if 'page' in request.GET:
        page = request.GET['page']
        offset = (int(request.GET['page']) - 1) * LIMIT
    # allcnt = utils.getResultBySql('select count(*) from account_user_tag ' + where)[0][0]

    try:
        allcntQuerySet = utils.SqlUtils().getArrResultWrapper('select count(*) from account_user_tag ' + where,logger,
                                                                                    'dashboardviews.py',
                                                                                    'getDataPerTag')
        allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0

        # fetchall = utils.getResultBySql(
        #     'select id,name,remark from account_user_tag ' +
        #     where + '  order by id desc LIMIT ' + str(LIMIT) + ' offset ' +
        #     str(offset))
        fetchall = utils.SqlUtils().getArrResultWrapper('select id,name,remark from account_user_tag ' +
                                                            where + '  order by id desc LIMIT ' + str(LIMIT) + ' offset ' +
                                                            str(offset),
                                                            logger,
                                                            'dashboardviews.py',
                                                            'getDataPerTag')
    except Exception as e:
        resultObj = tools.errorMes(e.args)
        return Response(resultObj)
    groups = []
    for obj in fetchall:
        dics = {}
        dics['id'] = obj[0]
        dics['tagname'] = obj[1]
        dics['tagremark'] = obj[2]
        groups.append(dics)
    resultObj = tools.successMes()
    resultObj['total'] = allcnt
    resultObj['rows'] = groups
    return Response(resultObj)


@api_view(http_method_names=['POST'])
# 用于分享类邮件
def shareTheme(request):
    shareType = request.data['shareType']
    users = request.data['users']
    picUrl = request.data['picUrl']
    filePath = 'code/datax/frontend/upload/temp_files/'
    if not os.path.exists(filePath):
        os.makedirs(filePath)
    ss = picUrl.split('base64,')
    ori_image_data = base64.b64decode(ss[1])
    imgPath = filePath + str(round(time.time() * 1000)) + '.png'
    with open(imgPath, 'wb') as fp:
        fp.write(ori_image_data)
    # 如果用户配置了邮件模板使用已配置的模板
    emailobj = emailconf.objects.filter()
    # 获取email配置信息。系统分享类模板所使用的模板id
    templateid = json.loads(emailobj[0].options)['emailsharetemplate']
    templateobj = msgalltemplateconf.objects.filter(id=templateid)
    if templateobj:
        mail_title = templateobj[0].templatetitle
        content = templateobj[0].templatecontent
    else:
        # 使用默认
        mail_title = '本次主题分享'
        content = '本次主题分享'
    mailTo = ''
    for user in users:
        if user['email'] is not None and len(user['email']) > 2:
            mailTo = user['email'] + ',' + mailTo
    if len(mailTo) > 3:
        pos = mailTo.rfind(',')
        mailTo = mailTo[:pos] + mailTo[pos + 1:]
    send_success = sendPicMail(mail_title, content, imgPath, mailTo, users)
    return Response({
        "status": send_success
    })


@api_view(http_method_names=['POST'])
def exportTheme(request):
    try:
        picUrl = request.data['picUrl']
        picWidth = request.data['picWidth']
        picHeight = request.data['picHeight']
        backgroundColor = request.data['backgroundColor']
        filePath = './frontend/upload/temp_files/'
        if not os.path.exists(filePath):
            os.makedirs(filePath)
        if backgroundColor == None or backgroundColor == 'null':
            backgroundColor = '#fff'
        fileType = request.data['fileType']
        # print(picUrl)
        print('picWidth:' + str(picWidth))
        print('picHeight:' + str(picHeight))
        print('fileType:' + fileType)
        ss = picUrl.split('base64,')
        ori_image_data = base64.b64decode(ss[1])
        imgPath = filePath + str(round(time.time() * 1000)) + '.png'
        with open(imgPath, 'wb') as fp:
            fp.write(ori_image_data)
        png2 = Image.open(imgPath)
        png2.load()
        if backgroundColor == 'noNeed':
            background = Image.new("RGB", png2.size, 'rgb(227, 230, 235)')
        else:
            background = Image.new("RGB", png2.size, backgroundColor)
        background.paste(png2, mask=png2.split()[3])  # 3 is the alpha channel
        if glvr.PIC_MARK_TEXT:
            # 设置所使用的字体
            font = ImageFont.truetype(glvr.PRJ_PATH + '/frontend/js/font/ziti.ttf', 24)

            # 画图
            draw = ImageDraw.Draw(background)
            draw.text((10, 10), glvr.PIC_MARK_TEXT, (255, 0, 0), font=font)  # 设置文字位置/内容/颜色/字体
            draw = ImageDraw.Draw(background)
        jpgName = str(round(time.time() * 1000)) + '.jpg'
        background.save(filePath + jpgName, 'JPEG', quality=120)

        if fileType == 'pic':
            fileName = jpgName

        if fileType == 'excel':
            fileName = str(round(time.time() * 1000)) + '.xlsx'
            workbook = xlsxwriter.Workbook(filePath + fileName)
            worksheet = workbook.add_worksheet()
            # worksheet.insert_image('D4', './frontend/upload/temp_files/test.jpg', {'x_scale': 0.9})
            worksheet.insert_image('B2', '', {'image_data': io.BytesIO(ori_image_data), 'x_scale': 0.9, 'y_scale': 1.1})
            workbook.close()
        elif fileType == 'word':
            fileName = str(round(time.time() * 1000)) + '.docx'
            document = Document()
            document.add_picture(io.BytesIO(ori_image_data), width=Inches(6))
            document.save(filePath + fileName)
        elif fileType == 'powerpoint':
            fileName = str(round(time.time() * 1000)) + '.pptx'
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            left = pptInches(0)
            top = pptInches(1)
            pic = slide.shapes.add_picture(filePath + jpgName, left, top, width=pptInches(10))
            prs.save(filePath + fileName)
        elif fileType == 'pdf':
            fileName = str(round(time.time() * 1000)) + '.pdf'
            (w, h) = landscape(A4)
            h = h + picHeight  ##加上图片本身的尺寸生成的pdf就比较规范
            w = w + picWidth
            # c = canvas.Canvas(filePath + fileName, pagesize = landscape(A4))
            c = canvas.Canvas(filePath + fileName, pagesize=(w, h))
            # c.drawString(100, 100, "Hello,World")

            c.drawImage(filePath + jpgName, 0, 0, w, h)
            # c.showPage()
            c.save()
    except Exception as e:
        raise
    # w = xlwt.Workbook()
    # ws = w.add_sheet('Image')
    # ws.write(1, 0, label='this is test')
    # ws.insert_bitmap("D:/FileTest/theme.bmp", 2, 2)

    # w.save(filePath)
    # response = HttpResponse(content_type='application/ms-excel;charset=GBK')
    # response['Content-Disposition'] = 'attachment; filename=beifen.xls'
    # w.save(response)
    return Response({
        "status": "ok",
        "filePath": filePath + fileName,
        "fileName": fileName
    })


@api_view(http_method_names=['POST'])
def setcharts(request):
    chartObj = request.data
    id = chartObj['id']  # 编辑对应id，如果新添加的则id=0 version = request.POST.get('no')
    name = chartObj['name']
    imgpath = chartObj['imgpath']
    filepath = ""
    returnObj = tools.successMes();
    if imgpath != '':
        t = time.time()
        year = datetime.datetime.now().year
        month = datetime.datetime.now().month
        day = datetime.datetime.now().day
        path = './frontend/upload/chart_img/' + str(year) + '-' + str(month) + '-' + str(day) + '/'
        filename = name + '_' + str(int(t)) + '.png'
        if os.path.exists(path):
            pass
        else:
            os.makedirs(path)
        try:
            filepath = path + filename
            if os.path.exists(filepath) and os.path.isfile(filepath):
                os.remove(filepath)
            ss = imgpath.split('base64,')
            # missing_padding = 4 - len(imgpath) % 4
            # if missing_padding:
            #     imgpath += b'=' * missing_padding
            imgdata = base64.b64decode(ss[1])
            with open(filepath, 'wb') as fp:
                fp.write(imgdata)
            img = Image.open(filepath)
            w, h = img.size
            if w > 500:
                h = int(round((500 / w) * h))
                w = 500
            img.resize((w, h)).save(filepath)
            filepath = '.' + filepath
        except Exception as e:
            filepath = ""
            a = e.args
            return Response(tools.errorMes(e.args))
    try:
        if not id:  # 添加
            row = charts.objects.create(name=name, kind=chartObj['kind'], echartconfig=chartObj['echartconfig'],
                                        charttype=chartObj['charttype'], jsonconfig=chartObj['jsonconfig'],
                                        createname=request.user, filterstring=chartObj['filterstring'],
                                        datasetstring=chartObj['datasetstring'],
                                        remark=chartObj['remark'], keywords=chartObj['keywords'],
                                        refreshspeed=chartObj['refreshspeed'], imgpath=filepath)
            returnObj['data']['id'] = row.id;
            # echartconfig data的url中添加chartid，并更换所使用的方法
            chartObj = charts.objects.get(id=row.id)
            echartconfig = json.loads(chartObj.echartconfig)
            echartconfig['data']['url'] = echartconfig['data']['url'].replace("getOlapData",
                                                                              "getOlapDataShow") + "/" + row.id
            chartObj.echartconfig = json.dumps(echartconfig)
            #  jsonconfig 的url中添加chartid，并更换所使用的方法
            datasetstring = json.loads(chartObj.datasetstring)
            #datasetstring['url'] = datasetstring['url'].replace("getOlapData", "getOlapDataShow") + "/" + row.id
            chartObj.datasetstring = json.dumps(datasetstring)
            chartObj.save()
            return Response(returnObj)
        else:  # 编辑
            charts.objects.filter(id=id).update(name=name, kind=chartObj['kind'],
                                                echartconfig=chartObj['echartconfig'],
                                                charttype=chartObj['charttype'],
                                                jsonconfig=chartObj['jsonconfig'],
                                                filterstring=chartObj['filterstring'],
                                                datasetstring=chartObj['datasetstring'],
                                                remark=chartObj['remark'], keywords=chartObj['keywords'],
                                                refreshspeed=chartObj['refreshspeed'], imgpath=filepath)
            returnObj['data']['id'] = id;
            #  解决之前组件显示问题，重新点击编辑保存后即可使用新的方法显示
            # echartconfig data的url中添加chartid，并更换所使用的方法
            chartObj = charts.objects.get(id=id)
            echartconfig = json.loads(chartObj.echartconfig)
            if echartconfig['data']['url'].find(id) == -1:
                echartconfig['data']['url'] = echartconfig['data']['url'].replace("getOlapData",
                                                                                  "getOlapDataShow") + "/" + id
                chartObj.echartconfig = json.dumps(echartconfig)
                chartObj.save()
            #  jsonconfig 的url中添加chartid，并更换所使用的方法
            datasetstring = json.loads(chartObj.datasetstring)
            # if datasetstring['url'].find(id) == -1:
            #     datasetstring['url'] = datasetstring['url'].replace("getOlapData", "getOlapDataShow") + "/" + id
            #     chartObj.datasetstring = json.dumps(datasetstring)
            #     chartObj.save()
            return Response(returnObj)
    except Exception as e2:
        print(e2)
        return Response(tools.errorMes(e2.args))


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getcharts(request):
    start = request.GET.get("startIndex", "0");  # &startIndex=0&count=10&orderBy=a.id+ASC 控件分页信息
    pagesize = request.GET.get("count", "5");
    name = request.GET.get("param[name]");
    dtype = request.GET.get("dtype");

    where = " where 1=1 ";
    if name:
        where += " and (name like '%" + name + "%' or keywords like '%" + name + "%') ";
    if dtype:
        where += " and kind in (" + dtype + ") ";
    # allcnt = utils.getResultBySql('select count(*) from dashboard_charts' + where)[0][0];
    # fetchall = utils.getResultBySql(
    #     "select a.id,a.name,a.kind,a.jsonconfig,a.filterstring,a.datasetstring,b.type_name,a.charttype,a.echartconfig,a.remark,a.refreshspeed,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname,a.keywords,a.allowusers, a.imgpath from dashboard_charts a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.createtime desc LIMIT " + pagesize + " offset " + start)
    fetchall = []
    executeSession = utils.SqlUtils()
    try:
        allcntQuerySet = executeSession.getArrResult('select count(*) from dashboard_charts' + where)
        fetchall = executeSession.getArrResult(
            "select a.id,a.name,a.kind,a.jsonconfig,a.filterstring,a.datasetstring,b.type_name,a.charttype,a.echartconfig,a.remark,a.refreshspeed,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname,a.keywords,a.allowusers, a.imgpath from dashboard_charts a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.createtime desc LIMIT " + pagesize + " offset " + start)
    except Exception as error:
        executeSession.rollBack()
        logger.error('---error---file:dashboardviews.py;method:getcharts;error:%s' % error)
    executeSession.closeConnect()
    allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0


    groups = []
    for obj in fetchall:
        dics = {}
        dics['id'] = obj[0]
        dics['name'] = obj[1]
        dics['kind'] = obj[6]
        dics['jsonconfig'] = obj[3]
        dics['filterstring'] = obj[4]
        dics['datasetstring'] = obj[5]
        dics['charttype'] = obj[7]
        dics['echartconfig'] = obj[8]
        dics['remark'] = obj[9]
        dics['refreshspeed'] = obj[10]
        dics['createtime'] = obj[11]
        dics['createname'] = obj[12]
        dics['keywords'] = obj[13]
        dics['allowusers'] = obj[14]
        dics['imgpath'] = obj[15]
        dics['datatype'] = 'chart'
        groups.append(dics)

    return Response({'totalcount': allcnt, 'datas': groups})


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getChartsTypelists(request):
    exesql = """select a.kind,b.type_name from dashboard_charts a left join dashboard_charttype b on a.kind = b.id GROUP BY a.kind,b.type_name"""
    # olaptypes = utils.getResultBySql(exesql)
    olaptypes = utils.SqlUtils().getArrResultWrapper(exesql,logger,'dashboardviews.py','getChartsTypelists')
    types = []
    for tp in olaptypes:
        if tp[0]:  # id不能为空
            type = {}
            type['id'] = tp[0]
            type['name'] = tp[1]
            types.append(type)
    return Response({'types': types})


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getAllChartsbgConf(request):
    exesql = """select"""
    # allChtbgConf=utils.getResultBySql(exesql)
    allChtbgConf = chartsbgconfig.objects.filter(status=0)
    allDt = []
    for chtbg in allChtbgConf:
        dt = {}
        dt['fileurl'] = chtbg.fileurl
        dt['confname'] = chtbg.confname
        allDt.append(dt)
    return Response({"data": allDt})


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getAllcharts(request):
    start = request.GET.get("startIndex", "0");  # &startIndex=0&count=10&orderBy=a.id+ASC 控件分页信息
    pagesize = request.GET.get("count", "5");
    name = request.GET.get("param[name]");
    where = " where 1=1 ";
    if name:
        where += " and name like '%" + name + "%' or keywords like '%" + name + "%' ";
    # fetchall = utils.getResultBySql(
    #     "select a.id,a.name,a.kind,a.jsonconfig,a.filterstring,a.datasetstring,b.type_name,a.charttype,a.echartconfig,a.remark,a.refreshspeed,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname,a.keywords,a.allowusers, a.imgpath from dashboard_charts a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.id desc")
    fetchall = utils.SqlUtils().getArrResultWrapper("select a.id,a.name,a.kind,a.jsonconfig,a.filterstring,a.datasetstring,b.type_name,a.charttype,a.echartconfig,a.remark,a.refreshspeed,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname,a.keywords,a.allowusers, a.imgpath from dashboard_charts a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.id desc",
                                                    logger,'dashboardviews.py', 'getAllcharts')
    groups = []
    for obj in fetchall:
        dics = {}
        dics['id'] = obj[0]
        dics['name'] = obj[1]
        dics['kind'] = obj[6]
        dics['jsonconfig'] = obj[3]
        dics['filterstring'] = obj[4]
        dics['datasetstring'] = obj[5]
        dics['charttype'] = obj[7]
        dics['echartconfig'] = obj[8]
        dics['remark'] = obj[9]
        dics['refreshspeed'] = obj[10]
        dics['createtime'] = obj[11]
        dics['createname'] = obj[12]
        dics['keywords'] = obj[13]
        dics['allowusers'] = obj[14]
        dics['imgpath'] = obj[15]
        dics['datatype'] = 'chart'
        groups.append(dics)

    # tablefetchall = utils.getResultBySql(
    #     "SELECT a.name,a.id,a.jsonconfig,b.type_name,a.createtime,a.createname,a.imgpath from dashboard_datatable a left join dashboard_charttype  b on a.kind = b.id ")
    tablefetchall = utils.SqlUtils().getArrResultWrapper(
        "SELECT a.name,a.id,a.jsonconfig,b.type_name,a.createtime,a.createname,a.imgpath from dashboard_datatable a left join dashboard_charttype  b on a.kind = b.id ",
        logger,'dashboardviews.py', 'getAllcharts')
    for obj in tablefetchall:
        dics = {}
        dics['id'] = obj[1]
        dics['name'] = obj[0]
        dics['jsonconfig'] = json.loads(obj[2].replace("'", "\""))
        dics['kind'] = obj[3]
        dics['filterstring'] = ""
        dics['datasetstring'] = ""
        dics['charttype'] = "table"
        dics['echartconfig'] = ""
        dics['remark'] = ""
        dics['refreshspeed'] = ""
        dics['createtime'] = obj[4]
        dics['createname'] = obj[5]
        dics['keywords'] = ''
        dics['allowusers'] = ''
        dics['imgpath'] = obj[6]
        dics['datatype'] = 'table'
        groups.append(dics)
    return Response({'datas': groups})


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getchartsper(request):
    start = request.GET.get("startIndex", "0");  # &startIndex=0&count=10&orderBy=a.id+ASC 控件分页信息
    pagesize = request.GET.get("count", "5");
    name = request.GET.get("param[name]");
    where = " where allowusers is null or allowusers='' or allowusers='[]' or allowusers like '%" + request.user.username + "%'  ";
    if name:
        where += " and name like '%" + name + "%'";
    # allcnt = utils.getResultBySql('select count(*) from dashboard_charts' + where)[0][0];
    # fetchall = utils.getResultBySql(
    #     "select a.id,a.name,a.kind,a.jsonconfig,a.filterstring,a.datasetstring,b.type_name,a.charttype,a.echartconfig,a.remark,a.refreshspeed,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname,a.keywords,a.allowusers from dashboard_charts a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.id desc LIMIT " + pagesize + " offset " + start)
    fetchall = []
    executeSession = utils.SqlUtils()
    try:
        allcntQuerySet = executeSession.getArrResult('select count(*) from dashboard_charts' + where)
        fetchall = executeSession.getArrResult("select a.id,a.name,a.kind,a.jsonconfig,a.filterstring,a.datasetstring,b.type_name,a.charttype,a.echartconfig,a.remark,a.refreshspeed,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname,a.keywords,a.allowusers from dashboard_charts a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.id desc LIMIT " + pagesize + " offset " + start)
    except Exception as error:
        executeSession.rollBack()
        logger.error('---error---file:dashboardviews.py;method:getchartsper;error:%s' % error)
    executeSession.closeConnect()
    allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0

    groups = []
    for obj in fetchall:
        dics = {}
        dics['id'] = obj[0]
        dics['name'] = obj[1]
        dics['kind'] = obj[6]
        dics['jsonconfig'] = obj[3]
        dics['filterstring'] = obj[4]
        dics['datasetstring'] = obj[5]
        dics['charttype'] = obj[7]
        dics['echartconfig'] = obj[8]
        dics['remark'] = obj[9]
        dics['refreshspeed'] = obj[10]
        dics['createtime'] = obj[11]
        dics['createname'] = obj[12]
        dics['keywords'] = obj[13]
        dics['allowusers'] = obj[14]
        groups.append(dics)
    return Response({'totalcount': allcnt, 'datas': groups})


@api_view(http_method_names=['GET'])
def delcharts(request):
    charts.objects.get(id=request.GET.get('id')).delete()
    return Response({'status': 'ok'})


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getchartRow(request, pk):
    row = charts.objects.get(id=pk);
    dist = {}
    dist['name'] = row.name
    dist['jsonconfig'] = row.jsonconfig
    dist['kind'] = row.kind
    dist['datastring'] = row.datasetstring
    dist['filterstring'] = row.filterstring
    return Response(dist)


@api_view(http_method_names=['POST'])
def setscenes(request):
    id = request.POST.get('id')  # 编辑对应id，如果新添加的则id=0 version = request.POST.get('no')
    name = request.POST.get('name')
    if not id:  # 添加
        row = scenes.objects.create(name=name, kind=request.POST.get("kind"),
                                    htmlrawstring=request.POST.get("htmlrawstring"),
                                    createname=request.user, htmlcleanconfig=request.POST.get("htmlcleanconfig")
                                    , remark=request.POST.get("remark"), keywords=request.POST.get("keywords"),
                                    basicconfig=request.POST.get("basicconfig"))
        return Response({
            "id": row.id,
            "status": "ok"
        })
    else:  # 编辑
        scenes.objects.filter(id=id).update(name=name, kind=request.POST.get("kind"),
                                            htmlrawstring=request.POST.get("htmlrawstring"),
                                            htmlcleanconfig=request.POST.get("htmlcleanconfig")
                                            , remark=request.POST.get("remark"), keywords=request.POST.get("keywords"),
                                            basicconfig=request.POST.get("basicconfig"))
        return Response({
            "id": id,
            "status": "ok"
        })


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getsenes(request):
    start = request.GET.get("startIndex", "0")  # &startIndex=0&count=10&orderBy=a.id+ASC 控件分页信息
    pagesize = request.GET.get("count", "5")
    name = request.GET.get("param[name]")
    dtype = request.GET.get("dtype")
    dashboardVersion = request.GET.get("version", "0")
    where = " where 1=1 "
    if name:
        where += " and (a.name like '%" + name + "%' or a.keywords like '%" + name + "%') "
    if dtype:
        where += " and a.kind in (" + dtype + ") "
    if str(dashboardVersion) == '1':
        where += " and (a.sceneversion < 2 OR a.sceneversion IS NULL) "
    # print('where=',where)
    # allcnt = utils.getResultBySql('select count(*) from dashboard_scenes a ' + where)[0][0]
    # fetchall = utils.getResultBySql(
    #     "select a.id,a.name,b.type_name,a.htmlrawstring,a.htmlcleanconfig,a.remark,a.keywords,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname from dashboard_scenes a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.createtime desc LIMIT " + pagesize + " offset " + start)
    fetchall = []
    executeSession = utils.SqlUtils()
    try:
        allcntQuerySet = executeSession.getArrResult('select count(*) from dashboard_scenes a ' + where)
        fetchall = executeSession.getArrResult(
            "select a.id,a.name,b.type_name,a.htmlrawstring,a.htmlcleanconfig,a.remark,a.keywords,to_char(a.createtime,'yyyy-mm-dd hh24:mi:ss')   createtime,a.createname from dashboard_scenes a left join dashboard_charttype b on a.kind = b.id " + where + " order by a.createtime desc LIMIT " + pagesize + " offset " + start)
    except Exception as error:
        executeSession.rollBack()
        logger.error('---error---file:dashboardviews.py;method:getsenes;error:%s' % error)
    executeSession.closeConnect()
    allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0

    groups = []
    for obj in fetchall:
        dics = {}
        dics['id'] = obj[0]
        dics['name'] = obj[1]
        dics['kind'] = obj[2]
        dics['htmlrawstring'] = obj[3]
        dics['htmlcleanconfig'] = obj[4]
        dics['remark'] = obj[5]
        dics['keywords'] = obj[6]
        dics['createtime'] = obj[7]
        dics['createname'] = obj[8]
        groups.append(dics)
    return Response({'totalcount': allcnt, 'datas': groups})


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getThemeTypelists(request):
    exesql = """select a.kind,b.type_name from dashboard_scenes a left join dashboard_charttype b on a.kind = b.id GROUP BY a.kind,b.type_name"""
    # olaptypes = utils.getResultBySql(exesql)
    olaptypes = utils.SqlUtils().getArrResultWrapper(exesql,logger,'dashboardviews.py','getThemeTypelists')
    types = []
    for tp in olaptypes:
        if tp[0]:  # id不能为空
            type = {}
            type['id'] = tp[0]
            type['name'] = tp[1]
            types.append(type)
    return Response({'types': types})


@api_view(http_method_names=['GET'])
def delscenes(request):
    scenes.objects.get(id=request.GET.get('id')).delete()
    return Response({'status': 'ok'})


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getsenceRow(request, id):
    row = scenes.objects.get(id=id);
    dist = {}
    dist['name'] = row.name
    dist['kind'] = row.kind
    dist['htmlrawstring'] = row.htmlrawstring
    dist['htmlcleanconfig'] = row.htmlcleanconfig
    dist['remark'] = row.remark
    dist['keywords'] = row.keywords
    dist['createtime'] = row.createtime.strftime("%Y-%m-%d %H:%M:%S")
    dist['createname'] = row.createname
    dist['basicconfig'] = row.basicconfig
    return Response(dist)


@api_view(http_method_names=['POST'])
def setthemes(request):
    id = request.POST.get('id')  # 编辑对应id，如果新添加的则id=0 version = request.POST.get('no')
    name = request.POST.get('name')
    if 'updateType' in request.POST:
        updateType = request.POST.get('updateType')
        if updateType == 'updateInterval' and len(id) > 0:
            theme = themes.objects.get(id=int(id))
            theme.themeconfig = request.POST.get("themeconfig")
            theme.save()
            return Response({
                "status": "ok"
            })
    if not id:  # 添加
        themes.objects.create(name=name, code=request.POST.get("code"), remark=request.POST.get("remark"),
                              scenesconfig=request.POST.get("scenesconfig"), createname=request.user,
                              themeconfig=request.POST.get("themeconfig"))
        return Response({
            "status": "ok"
        })
    else:  # 编辑
        themes.objects.filter(id=id).update(name=name, code=request.POST.get("code"), remark=request.POST.get("remark"),
                                            scenesconfig=request.POST.get("scenesconfig"),
                                            themeconfig=request.POST.get("themeconfig"))
        return Response({
            "status": "ok"
        })


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getthemes(request):
    start = request.GET.get("startIndex", "0");  # &startIndex=0&count=10&orderBy=a.id+ASC 控件分页信息
    pagesize = request.GET.get("count", "5");
    name = request.GET.get("param[name]");
    where = " where 1=1 ";
    if name:
        where += " and name like '%" + name + "%'";
    # allcnt = utils.getResultBySql('select count(*) from dashboard_scenes' + where)[0][0];
    # fetchall = utils.getResultBySql(
    #     'select id,name,code,scenesconfig,remark,themeconfig from dashboard_themes ' + where + ' order by createtime desc LIMIT ' + pagesize + ' offset ' + start)
    fetchall = []
    executeSession = utils.SqlUtils()
    try:
        allcntQuerySet = executeSession.getArrResult('select count(*) from dashboard_scenes' + where)
        fetchall = executeSession.getArrResult(
            'select id,name,code,scenesconfig,remark,themeconfig from dashboard_themes ' + where + ' order by createtime desc LIMIT ' + pagesize + ' offset ' + start)
    except Exception as error:
        executeSession.rollBack()
        logger.error('---error---file:dashboardviews.py;method:getthemes;error:%s' % error)
    executeSession.closeConnect()
    allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0

    groups = []
    for obj in fetchall:
        dics = {}
        dics['id'] = obj[0]
        dics['name'] = obj[1]
        dics['code'] = obj[2]
        dics['scenesconfig'] = obj[3]
        dics['remark'] = obj[4]
        dics['themeconfig'] = obj[5]
        groups.append(dics)
    return Response({'totalcount': allcnt, 'datas': groups})


@api_view(http_method_names=['GET'])
def delthemes(request):
    themes.objects.get(id=request.GET.get('id')).delete()
    return Response({'status': 'ok'})


@api_view(http_method_names=['GET'])
def getThemeDetail(request, pk, type):
    if type == 'scene':
        result = {}
        try:
            scene = scenes.objects.get(id=pk)
            lists = []
            dist = {}
            dist['id'] = pk
            dist['html'] = scene.htmlcleanconfig
            lists.append(dist)
            result['lists'] = lists
            result['name'] = scene.name
            result['basicconfig'] = scene.basicconfig
            try:
                result['options'] = json.loads(scene.options) if scene.options else ''
            except Exception as ee:
                result['options'] = ''

            result['code'] = '1'
        except Exception as e:
            result['code'] = '0'
        return Response(result)
    elif type == 'theme':
        dist = {}
        try:
            theme = themes.objects.get(id=pk)
            lists = json.loads(theme.scenesconfig.replace("'", "\""))
            for row in lists:
                id = row['id']
                scens = scenes.objects.get(id=id)
                row['html'] = scens.htmlcleanconfig
                row['basicconfig'] = scens.basicconfig
            dist['lists'] = lists
            dist['name'] = theme.name
            dist['themeconfig'] = theme.themeconfig
            dist['code'] = '1'
        except Exception as  e:
            dist['code'] = '0'
        return Response(dist)


@api_view(http_method_names=['POST'])
def setchartusers(request):
    id = request.POST.get('id')  # 编辑对应id，如果新添加的则id=0 version = request.POST.get('no')
    allowusers = request.POST.get('allowusers')
    charts.objects.filter(id=id).update(allowusers=allowusers)
    return Response({
        "status": "ok"
    })


@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getDataTable(request):
    where = """where 1 = 1 """
    if 'search' in request.GET:
        where = where + """ and a.name like '%""" + request.GET['search'] + """%' """
    if 'page' in request.GET:
        offset = (int(request.GET['page']) - 1) * 10
        # rows = utils.getResultBySql(
        #     """select a.id, a."name", a."remark",b.type_name, a.jsonconfig from dashboard_datatable a left join dashboard_charttype b on a.kind = b.id  """ + where + """   order by a.createtime desc LIMIT """
        #     + str(10) + ' offset ' + str(offset))
        sql = """select a.id, a."name", a."remark",b.type_name, a.jsonconfig, a."version" from dashboard_datatable a left join dashboard_charttype b on a.kind = b.id  """ + where + """   order by a.createtime desc LIMIT """\
            + str(10) + ' offset ' + str(offset)
        rows = utils.SqlUtils().getArrResultWrapper(sql,logger,'dashboardviews.py','getDataTable')

    else:
        # rows = utils.getResultBySql(
        #     """select a.id, a."name", a."remark",b.type_name, a.jsonconfig from dashboard_datatable a left join dashboard_charttype b on a.kind = b.id  """ + where + """   order by a.createtime desc"""
        # )
        sql = """select a.id, a."name", a."remark",b.type_name, a.jsonconfig, a."version" from dashboard_datatable a left join dashboard_charttype b on a.kind = b.id  """ + where + """   order by a.createtime desc"""
        rows = utils.SqlUtils().getArrResultWrapper(sql,logger, 'dashboardviews.py', 'getDataTable')

    # allcnt = utils.getResultBySql('select count(*) from dashboard_datatable a ' + where)[0][0]
    allcntQuerySet = utils.SqlUtils().getArrResultWrapper('select count(*) from dashboard_datatable a ' + where, logger,'dashboardviews.py', 'getDataTable')
    allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0

    groups = []
    for row in rows:
        dist = {}
        dist['name'] = row[1]
        dist['remark'] = row[2]
        dist['id'] = str(row[0])
        dist['kind'] = row[3]
        dist['json'] = row[4].replace("'", "\"")
        dist['version'] = row[5]
        groups.append(dist)
    return Response({'total': allcnt, 'rows': groups})


@api_view(http_method_names=['POST'])
def savedatatable(request):
    config = request.data['config']
    result = {}
    imgpath = config['imgpath']
    filepath = ""
    if imgpath != '':
        t = time.time()
        year = datetime.datetime.now().year
        month = datetime.datetime.now().month
        day = datetime.datetime.now().day
        path = './frontend/upload/chart_img/' + str(year) + '-' + str(month) + '-' + str(day) + '/'
        filename = config['title'] + '_' + str(int(t)) + '.png'
        if os.path.exists(path):
            pass
        else:
            os.makedirs(path)
        try:
            filepath = path + filename
            if os.path.exists(filepath) and os.path.isfile(filepath):
                os.remove(filepath)
            ss = imgpath.split('base64,')
            # missing_padding = 4 - len(imgpath) % 4
            # if missing_padding:
            #     imgpath += b'=' * missing_padding
            imgdata = base64.b64decode(ss[1])
            with open(filepath, 'wb') as fp:
                fp.write(imgdata)
            img = Image.open(filepath)
            w, h = img.size
            if w > 500:
                h = int(round((500 / w) * h))
                w = 500
            img.resize((w, h)).save(filepath)
            filepath = '.' + filepath
        except Exception as e:
            filepath = ""
            a = e.args
            pass
    if config['id'] == '':
        try:
            row = datatable.objects.get_or_create(name=config['title'], kind=config['kind'], jsonconfig=config['json'],
                                                  remark=config['desc'], imgpath=filepath)
            result['code'] = '1'
        except Exception as e:
            result['code'] = '0'
            result['msg'] = e.args
    else:
        row = datatable.objects.get(id=config['id'])
        row.name = config['title']
        row.remark = config['desc']
        row.kind = config['kind']
        row.jsonconfig = config['json']
        row.imgpath = filepath
        try:
            row.save()
            result['code'] = '1'
        except Exception as e:
            result['code'] = '0'
            result['msg'] = e.args
    return Response(result)


@api_view(http_method_names=['GET'])
def getDataTableRow(request, pk):
    row = datatable.objects.get(id=pk)
    dist = {}
    dist['name'] = row.name
    dist['remark'] = row.remark
    dist['kind'] = row.kind
    dist['json'] = json.loads(row.jsonconfig.replace("'", "\""))
    return Response(dist)


@api_view(http_method_names=['POST'])
def deleteDataTable(request):
    id = request.data['id']
    result = {}
    try:
        datatable.objects.filter(id=id).delete()
        result['code'] = '1'
    except Exception as e:
        result['code'] = '0'
        result['msg'] = e.args
    return Response(result)


@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def saveUploadFile(request):
    resultsObj = tools.successMes()
    try:
        f = request.FILES['filename']
        fileType = request.data['fileuptype']
        fileSuffix = f.name.split('.')[len(f.name.split('.'))-1]
        filename = fileType + '-' + time.strftime("%Y%m%d%H%M%S", time.localtime()) + str(random.randint(-10000, -100)) + '.' + fileSuffix
        if fileType == 'menuicon':
            path = glvr.PRJ_PATH + '/frontend/upload/menuicon/'
            if not os.path.exists(path):
                os.makedirs(path)
            with open(path + filename, 'wb+') as destination:
                for chunk in f.chunks():
                    destination.write(chunk)
            prjFilePath = (path.replace(glvr.PRJ_PATH, '')) + filename
            resultsObj['data']['filePath'] = prjFilePath
    except Exception as e:
        resultsObj = tools.errorMes(error.args)
    return Response(resultsObj)  

@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def uploadbgpicture(request):
    f = request.FILES['filename']
    filename = time.strftime("%Y%m%d%H%M%S", time.localtime()) + f.name
    id = request.user.id
    path = './frontend/upload/user_' + str(id) + '/'
    if os.path.exists(path):
        a = 1
    else:
        os.makedirs(path)
    try:
        distpath = path + filename
        if os.path.exists(distpath) and os.path.isfile(distpath):
            os.remove(distpath)
        with open(path + filename, 'wb+') as destination:
            for chunk in f.chunks():
                destination.write(chunk)
    except Exception as e:
        a = e.args
    return Response({'code': 1, "path": '/frontend/upload/user_' + str(id) + '/', "filename": filename})

@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def uploadchartsbgpicture(request):
    f = request.FILES['filename']
    filename = time.strftime("%Y%m%d%H%M%S", time.localtime()) + f.name
    id = request.user.id
    path = './frontend/upload/user_' + str(id) + '/chartsbgimgs/'
    if os.path.exists(path):
        a = 1
    else:
        os.makedirs(path)
    try:
        distpath = path + filename
        if os.path.exists(distpath) and os.path.isfile(distpath):
            os.remove(distpath)
        with open(path + filename, 'wb+') as destination:
            for chunk in f.chunks():
                destination.write(chunk)
    except Exception as e:
        a = e.args
    return Response({'code': 1, "path": '/frontend/upload/user_' + str(id) + '/chartsbgimgs/', "filename": filename})


@api_view(http_method_names=['POST'])
def uploadMultiplePic(request):
    results = tools.successMes()
    results['data'] = []
    count = 0
    path = './frontend/image/newdash/text_editor_imgs/'
    if os.path.exists(path):
        a = 1
    else:
        os.makedirs(path)
    for imgFile in request.FILES:
        f = request.FILES[imgFile]
        count += 1
        filename = time.strftime("%Y%m%d%H%M%S", time.localtime()) + str(count) + f.name
        try:
            distpath = path + filename
            if os.path.exists(distpath) and os.path.isfile(distpath):
                os.remove(distpath)
            with open(path + filename, 'wb+') as destination:
                for chunk in f.chunks():
                    destination.write(chunk)
            results['data'].append({'url': distpath[1:]})
        except Exception as e:
            print('file:getFilterDatas; method:uploadMultiplePic')
            print(e)
            return Response(tools.errorMes(e.args))

    return Response(results)


@api_view(http_method_names=['POST'])
def upload_file_conf(request):
    """ 上传文件，并将信息入库"""
    rs = {}
    f = request.FILES['filename']
    upload_file_name = f.name
    # 保存文件到/frontend/upload/temp_files
    save_path = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    write_file_name = datetime.datetime.now().strftime("%Y-%m-%d_%H_%M_%S")  # 写入到本地的文件名
    savePath = os.path.join(save_path, 'frontend', 'upload',
                            'packageConfig') + os.sep + write_file_name + upload_file_name
    print(savePath)
    with open(savePath, 'wb') as fw:  # 写入到指定目录
        fw.write(f.read())
    file_replace(savePath)
    # 保存记录
    licenseconf.objects.create(username=request.user.username, filename=upload_file_name)

    try:
        rs,
    except Exception as e:
        rs['code'] = 0
        rs['msg'] = e.args
        raise e
    return Response(rs)


def file_replace(fileUrl):
    """ 文件替换 """
    pathdir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    filePathdir = os.path.join(pathdir, 'datax')
    originalFilePath = filePathdir + '/datax_license.xml'
    print(originalFilePath)
    with open(fileUrl) as file_obj:
        conent = file_obj.readlines()
        for i in range(0, len(conent)):
            conent[i] = conent[i]
    with open(originalFilePath, 'w') as file_obj:
        file_obj.writelines(conent)


@api_view(http_method_names=['GET'])
def get_imp_licence_logs(request):
    """ 获取导入的licence信息 """
    where = ''
    if 'page' in request.GET:
        page = request.GET['page']
        offset = (int(page) - 1) * LIMIT
    # allcnt = utils.getResultBySql('select count(*) from dashboard_licenseconf' + where)
    # fetchall = utils.getResultBySql(
    #     'select filename,username,createdate from dashboard_licenseconf' +
    #     where + ' order by createdate desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset)
    # )
    fetchall = []
    executeSession = utils.SqlUtils()
    try:
        allcntQuerySet = executeSession.getArrResult('select count(*) from dashboard_licenseconf' + where)
        fetchall = executeSession.getArrResult(
            'select filename,username,createdate from dashboard_licenseconf' +
            where + ' order by createdate desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset))
    except Exception as error:
        executeSession.rollBack()
        logger.error('---error---file:dashboardviews.py;method:get_imp_licence_logs;error:%s' % error)
    executeSession.closeConnect()
    allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0

    allLicenseLogs = []
    for obj in fetchall:
        dt = {}
        dt['filename'] = obj[0]
        dt['username'] = obj[1]
        if obj[2]:
            dt['createdate'] = obj[2].strftime("%Y-%m-%d %H:%M:%S")
        else:
            dt['createdate'] = obj[2]

        allLicenseLogs.append(dt)

    return Response({'total': allcnt, 'rows': allLicenseLogs})


# region 获取场景列表
@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def getSceneList(request):
    results = tools.successMes()
    try:
        # 获取表数据
        sql = 'SELECT {0} FROM dashboard_scenes t1 LEFT JOIN dashboard_charttype t2 ON t1.kind = t2.id WHERE t1."sceneversion" >= 2 '
        colStr = '''
             t1."id", t1."name", t2."type_name", t1."remark", t1."keywords", t1."basicconfig", 
            t1."options", t1."modifytime"
        '''
        listSql = sql.format(colStr)
        countSql = sql.format('count(t1.*)')
        # 拼接查询条件
        where = " "
        if 'filterConfig' in request.data:
            filterConfig = request.data['filterConfig']
            if 'name' in filterConfig and len(filterConfig['name']) > 0:
                where = where + ' AND "name" LIKE \'%' + filterConfig['name'] + '%\' '
            if 'tags' in filterConfig and len(filterConfig['tags']) > 0:
                for tag in filterConfig['tags']:
                    val = tag
                    tempWhere = ' AND ({0})'
                    tempParams = '"keywords" LIKE \'' + val + ',%\' OR "keywords" LIKE \'%,' \
                                 + val + ',%\' OR "keywords" LIKE \'%,' + val + '\' OR "keywords" = \'' + val + '\'' \
                                 + ' OR "keywords" LIKE \'' + val + '，%\' OR "keywords" LIKE \'%，' \
                                 + val + '，%\' OR "keywords" LIKE \'%，' + val + '\' '
                    tempWhere = tempWhere.format(tempParams)
                    where = where + tempWhere
            listSql = listSql + where
        listSql = listSql + ' ORDER BY t1."modifytime" DESC'
        # 拼接分页条件
        if 'page' in request.data:
            if 'limit' in request.data:
                limit = request.data['limit']
            else:
                limit = LIMIT
            offset = (int(request.data['page']) - 1) * limit
            listSql = listSql + ' LIMIT ' + str(limit) + ' offset ' + str(offset)

        # dataTable = utils.getResultBySql(listSql)
        dataTable = utils.SqlUtils().getArrResultWrapper(listSql,logger,'dashboardviews.py', 'getSceneList')
        rows = []
        jsonStr='{}'
        for dataRow in dataTable:
            row = {}
            row['id'] = dataRow[0]
            row['name'] = dataRow[1]
            row['type_name'] = dataRow[2]
            row['remark'] = dataRow[3]
            row['keywords'] = dataRow[4]
            if dataRow[5] is not None:
                matchStr = '[\']{0,1}["]{0,1}imgUrl[\']{0,1}["]{0,1}:\s*[\']{0,1}["]{0,1}[\.]{0,1}/frontend/upload/chart_img/scene_imgs/[0-9\.]*png[\']{0,1}["]{0,1}'
                if re.search(matchStr, str(dataRow[5])):
                    m = re.search(matchStr, str(dataRow[5]))
                    jsonStr = '{' + (m.group(0)).replace("'", "\"") + '}'
                row['basicconfig'] = json.loads(jsonStr)
            else:
                row['basicconfig'] = {}
            row['options'] = dataRow[6]
            row['modifytime'] = dataRow[7]
            rows.append(row)
        # total = utils.getResultBySql(countSql)[0][0]
        totalQuerySet = utils.SqlUtils().getArrResultWrapper(countSql,logger,'dashboardviews.py', 'getSceneList')
        total = totalQuerySet[0][0] if totalQuerySet else 0

        results['rows'] = rows
        results['total'] = total
    except Exception as e:
        print('file:dashboardview; method:getSceneList')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)

# region 场景中的select过滤器取值
@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def getFilterDatas(request):
    results = tools.successMes()
    try:
        olapCols = request.data['olapCols']
        rows = []
        sql = 'SELECT DISTINCT col FROM ('
        for olapCol in olapCols:                #循环传入的参数，如果是数据提取olap就拼接sql最后执行unin all的sql，如果是实时或直连就先计算出结果再在最后拼接两个结果集
            if not olapCol['table'] or olapCol['table'].lower().find('none') > -1 or \
                olapCol['table'].lower().find('null') > -1 or olapCol['table'].lower().find('nan') > -1 :
                tempOlapObj = olap.objects.get(id=olapCol['id'])
                tempSourceObj = source.objects.get(id=tempOlapObj.sourceid)
                dirConnSql = ' SELECT DISTINCT ' + olapCol['selected'] + ' AS col FROM (' + tempSourceObj.sql + ')temptable1 where ' + olapCol['selected'] + ' is not null'
                #  晋能oa报表（部门考勤报表根据钉钉当前登陆人读取部门和子部门，只能够看到自己的信息）
                if olapCol['id'] == '593ed48ad43a11e9ba0d005056bac0e6' or olapCol['id'] == '28c0fda6d45d11e9857d005056bac0e6'or olapCol['id'] == 'cc13f512d8f411e9b6f3005056bac0e6':
                    departmentname = request.query_params['departmentname']
                    parentDepartmentname = request.query_params['parentDepartmentname']
                    if olapCol['selected'] == 'userdefinedsql0__superdep' :
                        # 该部门领导可以看到当前部门所有子部门信息
                        if parentDepartmentname == '集团':
                            dirConnSql = dirConnSql + " and userdefinedsql0__superdep='"+departmentname+"'"
                        else:                           
                            dirConnSql = dirConnSql + " and userdefinedsql0__superdep='"+parentDepartmentname+"'"
                    elif  olapCol['selected'] =='userdefinedsql0__departmentname':
                        if parentDepartmentname == '集团':
                            dirConnSql = dirConnSql + " and userdefinedsql0__superdep='"+departmentname+"'"
                        else:
                            dirConnSql = dirConnSql + " and userdefinedsql0__departmentname='"+departmentname+"'" 
                 #  晋能oa报表（个人审批分析报表根据钉钉当前登陆人读取信息，只能够看到自己的信息）                
                if olapCol['id'] == '677cb586e0f711e9b5d7005056bac0e6' or olapCol['id'] == '5945690ae0ea11e9b928005056bac0e6':
                    if 'username' in request.query_params:
                        username = request.query_params['username']
                        dirConnSql = dirConnSql + " and userdefinedsql0__lastname='"+username+"'"
                    
                #dirConnAdditionSql用于对select字段添加引号，有的数据库字段需要加引号
                dirConnAdditionSql = ' SELECT DISTINCT "' + olapCol['selected'] + '" AS col FROM (' + tempSourceObj.sql + ')temptable1 ' + olapCol['selected'] + ' is not null'
                st = stRestoreBySourceId(tempOlapObj.sourceid)
                try:
                  
                    resultRows = st.execute(dirConnSql).fetchall()
                except Exception as error:
                    print('--some exception raised,but i deal with it use anthor function,then the error=',error)
                    resultRows = st.execute(dirConnAdditionSql).fetchall()

                for resultRow in resultRows :  
                    rows.append(resultRow)
            else:
                sql = sql + ' SELECT "' + olapCol['selected'] + '" AS col FROM ' + olapCol['table'] + ' UNION ALL '
        olapConnRws = []
        if len(sql) > 26 :
            sql = sql[:-10] + ') t1 WHERE t1.col IS NOT NULL ORDER BY t1.col ASC'
            # olapConnRws = utils.getResultBySql(sql, utils.DATAXEXTENSION_DB_CHAR)
            olapConnRws = utils.SqlUtils(DBTypeCode.EXTENSION_DB.value).getArrResultWrapper(sql,logger,'dashboardviews.py', 'getDataTable')

        rows = rows + olapConnRws   #把两个结果集合并起来
        duplicationRows = []#对rows去重复再排序
        for row in rows :#去重复
            if row and str(row[0]).lower() not in ['','null','nan','none'] and str(row[0]) not in duplicationRows :
                duplicationRows.append(str(row[0]))
            # if row and row[0] and row[0] not in duplicationRows :
            #     duplicationRows.append(row[0])
            # elif row[0] == 0 :#0值必须保留
            #     duplicationRows.append(row[0])
        duplicationRows.sort()#排序

        dist = []
        for row in duplicationRows:
            obj = {}
            obj['id'] = row
            obj['name'] = row
            obj['checked'] = False
            obj['show'] = True
            dist.append(obj)
        results['data'] = dist
    except Exception as e:
        print('file:getFilterDatas; method:getFilterDatas')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


# endregion

# 获取场景中的所有olap字段
@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def getOlapTableCols(request):
    results = tools.successMes()
    try:
        olapIds = request.data['olapIds']
        # 获取所有olap
        olapObjs = olap.objects.filter(id__in=olapIds)
        olapcolumnObjs = olapcolumn.objects.filter(olapid__in=olapIds)
        olaps = utils.model_list_to_dict_wrapper(olapObjs)
        colsTemp = utils.model_list_to_dict_wrapper(olapcolumnObjs)
        colsObjs = []
        colDf = pd.DataFrame(colsTemp)
        colDfGroup = colDf.groupby(['olapid'])
        for olapId, group in colDfGroup:
            colsObj = {}
            colsObj['id'] = olapId
            colsObj['cols'] = group.to_dict(orient='records')
            colsObj['selected'] = ''
            colsObj['textFilterWay'] = 'equal'
            for olapObj in olaps:
                if olapId == olapObj['id']:
                    colsObj['name'] = olapObj['name']
                    colsObj['table'] = olapObj['table']
            colsObjs.append(colsObj)
        results['olaps'] = olaps
        results['colsObjs'] = colsObjs
    except Exception as e:
        print('file:dashboardviews; method:getOlapTableCols')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


# 保存场景设计
@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def saveSceneConfig(request):
    results = tools.successMes()
    try:
        sceneData = request.data['sceneObj']
        previewImgData = request.data['previewImgData']
        if len(previewImgData) > 0:
            filePath = glvr.PRJ_PATH + '/frontend/upload/chart_img/scene_imgs/'
            if not os.path.exists(filePath):
                os.makedirs(filePath)
            ss = previewImgData.split('base64,')
            ori_image_data = base64.b64decode(ss[1])
            imgPath = filePath + str(round(time.time() * 1000)) + '.png'
            with open(imgPath, 'wb') as fp:
                fp.write(ori_image_data)
            print('imagePath==', imgPath)
            srcImg = Image.open(imgPath)  # 修改图片大小
            srcImgSize = srcImg.size
            newSrcImgSize = [350, 350]
            if srcImgSize and srcImgSize[0] > 0 and srcImgSize[1] > 0:
                if srcImgSize[0] > srcImgSize[1]:
                    newSrcImgSize[1] = round(newSrcImgSize[0] * srcImgSize[1] / srcImgSize[0])
                else:
                    newSrcImgSize[0] = round(newSrcImgSize[1] * srcImgSize[0] / srcImgSize[1])
            newImg = srcImg.resize(tuple(newSrcImgSize), Image.ANTIALIAS)
            newImg.save(imgPath)  # 修改图片大小，完毕
            if 'imgUrl' in sceneData['basicconfig'] and sceneData['basicconfig']['imgUrl'] and len(sceneData['basicconfig']['imgUrl']) > 0 and os.path.exists(glvr.PRJ_PATH + sceneData['basicconfig']['imgUrl']):
                os.remove(glvr.PRJ_PATH + sceneData['basicconfig']['imgUrl'])
            sceneData['basicconfig']['imgUrl'] = imgPath.replace(glvr.PRJ_PATH, '')
            #给imgPath合并背景图，背景图来源于globalStyle的backgroundImage，如果没有背景图就设置为灰色
            globalBkImgPath = ''
            if 'globalStyle' in sceneData['options']:
                if 'backgroundImage' in sceneData['options']['globalStyle'] and \
                    sceneData['options']['globalStyle']['backgroundImage']:
                    globalBkImgPath = sceneData['options']['globalStyle']['backgroundImage']
            bkImg = ''
            # srcImg#截的图
            srcImg = Image.open(imgPath)  # 这个时候的图片已经调整了大小
            if globalBkImgPath:#如果有背景图片就合并背景图片
                globalBkImgPath = '.' + globalBkImgPath
                if os.path.exists(globalBkImgPath):
                    bkImg = Image.open(globalBkImgPath)
                    bkImg = bkImg.resize(newSrcImgSize, Image.ANTIALIAS)#背景图片调整大小为xxx
            if bkImg:
                srcImgR, srcImgG, srcImgB, srcImgRAlpha = srcImg.split()
                srcImgRAlpha = srcImgRAlpha.point(lambda i: i > 70 and 220)#i是初始像素，如果小于70就什么也不填写，大于70填255
                bkImg = Image.composite(srcImg, bkImg, srcImgRAlpha)
                bkImg.save(imgPath)#如果有bkImg则必定有globalBkImgPath，重写图片
            else:#用灰色铺满背景
                background = Image.new("RGB", srcImg.size, 'rgb(227, 230, 235)')
                background.paste(srcImg, mask=srcImg.split()[3])  # 3 is the alpha channel
                background.save(imgPath)

        configitems = sceneData['items']
        sceneOptions = json.dumps(sceneData['options']) if sceneData['options'] else None
        if len(sceneData['id']) > 1:
            sceneObj = scenes.objects.get(id=sceneData['id'])
            if sceneObj.basicconfig is not None:
                matchStr = '[\']{0,1}["]{0,1}imgUrl[\']{0,1}["]{0,1}:\s*[\']{0,1}["]{0,1}[\.]{0,1}/frontend/upload/chart_img/scene_imgs/[0-9\.]*png[\']{0,1}["]{0,1}'
                if re.search(matchStr, str(sceneObj.basicconfig)):
                    m = re.search(matchStr, str(sceneObj.basicconfig))
                    urlStr = m.group(0)
                    if urlStr[:1] is not '.':
                        urlStr = '.' + urlStr
                    if os.path.exists(urlStr):
                        os.remove(urlStr)
            sceneObj.items = json.dumps(configitems)
            sceneObj.name = sceneData['name']
            sceneObj.kind = sceneData['kind']
            sceneObj.remark = sceneData['remark']
            sceneObj.keywords = sceneData['keywords']
            sceneObj.basicconfig = json.dumps(sceneData['basicconfig'])
            sceneObj.sceneversion = glvr.DASHBOARD2_VERSION
            sceneObj.options = sceneOptions
            if sceneObj.createname is None or sceneObj.createname == '':
                sceneObj.createname = request.user.id
            sceneObj.save()

        else:

            sceneObj = scenes.objects.create(items=json.dumps(configitems), name=sceneData['name'],
                                             kind=sceneData['kind'], remark=sceneData['remark'],
                                             createname=request.user.id,
                                             keywords=sceneData['keywords'],
                                             basicconfig=json.dumps(sceneData['basicconfig']),
                                             sceneversion=glvr.DASHBOARD2_VERSION, options=sceneOptions)
        results['sceneObj'] = utils.model_to_dict_wrapper(sceneObj)
    except Exception as e:
        print('file:dashboardviews; method:saveSceneConfig')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


# 临时保存
@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def quickSave(request):
    results = tools.successMes()
    try:
        sceneId = request.data['sceneId']
        dataObj = json.loads(request.data['objStr'].replace("'", "\""))
        saveType = request.data['saveType']
        sceneObj = scenes.objects.get(id=sceneId)
        if saveType == 'singleItem':
            items = json.loads(sceneObj.items)
            isChange = False
            for item in items:
                if item['data']['id'] == dataObj['data']['id']:
                    for obj in dataObj['data']['commentsConfig']['comments']:
                        if obj['status'] != 'delete':
                            existFlag = False
                            for oldObj in item['data']['commentsConfig']['comments']:
                                if obj['id'] == oldObj['id']:
                                    oldObj['text'] = obj['text']
                                    existFlag = True
                                    break
                            if not existFlag:
                                item['data']['commentsConfig']['comments'].insert(0, obj)
                    isChange = True
                    results['item'] = item
                    break
            if not isChange:
                items.append(dataObj)
                results['item'] = dataObj
            sceneObj.items = json.dumps(items)
            sceneObj.save()
    except Exception as e:
        print('file:dashboardviews; method:quickSave')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


# 加载场景设计
@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def getNewScene(request):
    results = tools.successMes()
    try:
        sceneId = request.data['sceneId']
        sceneObj = scenes.objects.get(id=sceneId)
        results['sceneObj'] = utils.model_to_dict_wrapper(sceneObj)
    except Exception as error:
        logger.error('---error---file:dashboardviews; method:getScene; error=%s' % error)
        return Response(tools.errorMes(error.args))
    return Response(results)


# 在场景设计中导出样式包
@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def exportStyleFromScenesConfig(request):
    resultsObj = tools.successMes()
    try:
        scenesObj = request.data['scenesObj']
        gridsterMargins = request.data['gridsterMargins']
        osdir = glvr.PRJ_PATH
        timeStampSuff = str(int(time.time()))
        copyObj = {
            'confname': 'scenes_config_' + timeStampSuff,
            'configcontent': {
                'componentMargin': {'marginLeft': gridsterMargins[0], 'marginTop': gridsterMargins[1]},
                'globalStyle': scenesObj['options']['globalStyle'],
                'toolBarStyle': scenesObj['options']['toolBarStyle'],
                'filterStyle': scenesObj['options']['filterStyle'],
                'chartStyle': scenesObj['options']['chartStyle'],
                'syncFontColor': scenesObj['options']['syncFontColor'],
                'commentsStyle': scenesObj['options']['commentsStyle'],
            },
            'remark': '',
            'status': 0,
            'options': '',
            'customconfig': scenesObj['basicconfig'],
            'filecontent': '',
            'filename': '',
            'fileurl': '',
        }
        jsonConfig = copyObj['configcontent']
        globalBGImage = ''
        if 'globalStyle' in jsonConfig and 'backgroundImage' in jsonConfig['globalStyle']:
            globalBGImage = osdir + os.sep + jsonConfig['globalStyle']['backgroundImage']
        chartBGImage = ''
        if 'chartStyle' in jsonConfig and 'backgroundImage' in jsonConfig['chartStyle']:
            chartBGImage = osdir + os.sep + jsonConfig['chartStyle']['backgroundImage']
        filterBGImage = ''
        if 'filterStyle' in jsonConfig and 'backgroundImage' in jsonConfig['filterStyle']:
            filterBGImage = osdir + os.sep + jsonConfig['filterStyle']['backgroundImage']
        exportBGConfigParentDir = r'/frontend/upload/temp_files/sceneBGConfig_' + timeStampSuff + '/'
        exportBGConfigDir = osdir + exportBGConfigParentDir
        if os.path.exists(exportBGConfigDir):
            shutil.rmtree(exportBGConfigDir)#清空文件夹，但改文件夹也被清空了
        os.mkdir(exportBGConfigDir)#创建文件夹
        copyObj['configcontent'] = json.dumps(copyObj['configcontent'])
        copyObj['customconfig'] = json.dumps(copyObj['customconfig'])
        with open(exportBGConfigDir + 'config.json','w') as writer:#将json数据写入config.json
            json.dump(copyObj,writer)
        #将图片写入文件夹
        if globalBGImage :
            fileName = os.path.basename(globalBGImage)
            if fileName:#这里不判断上传文件的图片类型
                shutil.copy(globalBGImage,exportBGConfigDir + fileName)#将图片拷贝到exportBGConfigParentDir目录下
        if chartBGImage:
            fileName = os.path.basename(chartBGImage)
            if fileName:  # 这里不判断上传文件的图片类型
                shutil.copy(chartBGImage, exportBGConfigDir + fileName)  # 将图片拷贝到exportBGConfigParentDir目录下
        if filterBGImage:
            fileName = os.path.basename(filterBGImage)
            if fileName:  # 这里不判断上传文件的图片类型
                shutil.copy(filterBGImage, exportBGConfigDir + fileName)  # 将图片拷贝到exportBGConfigParentDir目录下
        #将文件压缩成zip文件并提供给js下载
        zipFileName = exportBGConfigDir[:-1] + '.zip'
        zf = zipfile.ZipFile(zipFileName,'w',zipfile.ZIP_DEFLATED)
        for dirPath,dirNames,fileNames in os.walk(exportBGConfigDir):
            # fPath = dirPath.replace(exportBGConfigDir,'')#需要替换根目录，不这样做就会从根目录复制
            fPath = 'sceneBGConfig_' + timeStampSuff + os.sep #需要替换根目录，不这样做就会从根目录复制
            for fileName in fileNames:
                zf.write(os.path.join(dirPath,fileName),fPath + fileName)
        logging.info('----场景配置数据包压缩成功-----')
        zf.close()
        resultsObj['downLoadZipFilePath'] = exportBGConfigParentDir[:-1] + '.zip'    #将路径返回给js以供下载
        resultsObj['timeStampSuff'] = timeStampSuff 
    except Exception as error:
        resultsObj = tools.errorMes(error.args)
        print(error)
        logger.error('---error---file:api.dashboardapi.dashboardview.py;method:exportStyleFromScenesConfig;error=%s' % error)
    return Response(resultsObj)


# region 主题、场景收藏


@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def getUserCollection(request):
    results = tools.successMes()
    try:
        collectType = request.data['collectType']
        currentuser = get_user(request)
        # collections = usercollection.objects.filter(user_id=currentuser.id, collect_type=collectType)
        # results['collections'] = utils.model_list_to_dict_wrapper(collections)
        sql = ''
        if collectType == 'scene':
            sql = '''
                SELECT t.*, t2.name AS name, t2.kind AS kind, t2.keywords AS tags, t2.sceneversion AS sceneversion, '-1' AS themeversion
                FROM dashboard_usercollection t
                LEFT JOIN dashboard_scenes t2 ON t.collect_id = t2.id
                WHERE t.collect_type = '{0}' and t.user_id='{1}'
            '''
        elif collectType == 'theme':
            sql = '''
                SELECT t1.*, t3.name AS name, t3.kind AS kind, t3.keywords AS tags, '-1' AS sceneversion, '-1' AS themeversion
                FROM dashboard_usercollection t1 
                LEFT JOIN dashboard_themes t3 ON t1.collect_id = t3.id
                WHERE t1.collect_type = '{0}' and t1.user_id='{1}' 
            '''
        else:
            sql = '''
                SELECT t.*, t2.name AS name, t2.kind AS kind, t2.keywords AS tags, t2.sceneversion AS sceneversion, '-1' AS themeversion
                FROM dashboard_usercollection t
                LEFT JOIN dashboard_scenes t2 ON t.collect_id = t2.id
                WHERE t.collect_type = '{0}' and t.user_id='{1}'
                UNION ALL
                SELECT t1.*, t3.name AS name, t3.kind AS kind, t3.keywords AS tags, '-1' AS sceneversion, '-1' AS themeversion
                FROM dashboard_usercollection t1 
                LEFT JOIN dashboard_themes t3 ON t1.collect_id = t3.id
                WHERE t1.collect_type = '{0}' and t1.user_id='{1}' 
            '''
        sql = sql.format(collectType, currentuser.id)
        # sqlResults = utils.getResultBySql(sql)
        sqlResults = utils.SqlUtils().getArrResultWrapper(sql,logger, 'dashboardviews.py','getUserCollection')
        results['collections'] = []
        for row in sqlResults:
            rowObj = {
                'id': row[0],
                'user_id': row[1],
                'collect_type': row[2],
                'collect_id': row[3],
                'options': row[4] if row[0] else {},
                'modify_time': row[5],
                'name': row[6],
                'kind': row[7],
                'tags': row[8],
                'sceneversion': row[9],
                'themeversion': row[10],
            }
            results['collections'].append(rowObj)
    except Exception as e:
        print('file:dashboardviews; method:getUserCollection')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


@api_view(http_method_names=['POST'])
@permission_classes((permissions.AllowAny,))
def saveUserCollection(request):
    results = tools.successMes()
    try:
        collectType = request.data['collectType']
        collectId = request.data['collectId']
        operation = request.data['operation']
        currentuser = get_user(request)
        if operation == 'connect':
            usercollection.objects.create(user_id=currentuser.id, collect_id=collectId,
                                          collect_type=collectType)
        else:
            usercollection.objects.filter(user_id=currentuser.id, collect_id=collectId,
                                          collect_type=collectType).delete()
    except Exception as e:
        print('file:dashboardviews; method:getUserCollection')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


# endregion

@api_view(http_method_names=['GET'])
def get_Unread_msg_info(request):
    """ 获取未读消息列表 """
    # 获取普通消息
    ordinary_msg = []
    where = 'mail_to=' + "'" + request.user.email + "' and is_read = '0'"
    print ('whereaaaa', where)
    # 条数
    sql = 'select count(*) from connect_maillogs where ' + where

    fetchall_ordinary_msg = []
    executeSession = utils.SqlUtils()
    try:
        msg_numberQuerySet = executeSession.getArrResult(sql)
        fetchall_ordinary_msg = executeSession.getArrResult(
            'select id,mail_title,mail_from,mail_content,create_date,img_path,is_read from connect_maillogs where' + ' ' + where)
    except Exception as error:
        executeSession.rollBack()
        logger.error('---error---file:dashboardviews.py;method:get_all_msg_info;error:%s' % error)
    executeSession.closeConnect()
    msg_number = msg_numberQuerySet[0][0] if msg_numberQuerySet else 0

    for obj in fetchall_ordinary_msg:
        dt = {}
        dt['id'] = obj[0]
        dt['mail_title'] = obj[1]
        dt['mail_from'] = obj[2]
        dt['mail_content'] = obj[3]
        dt['create_date'] = obj[4].strftime("%Y-%m-%d %H:%M:%S")
        dt['img_path'] = obj[5]
        dt['is_read'] = obj[6]
        ordinary_msg.append(dt)
    return Response({'msg_number': msg_number, 'ordinary_msg': ordinary_msg})


@api_view(http_method_names=['GET'])
def get_message_info(request):
    """获取消息详细内容信息"""
    if 'id' in request.GET:
        id = request.GET['id']
    # mark用来区分消息详情的展示（类别区分）
    if 'mark' in request.GET:
        mark = request.GET['mark']
        if mark == '0' or mark == '2':
            result = maillogs.objects.get(id=id)
            dt = {}
            dt['mail_title'] = result.mail_title
            dt['mail_from'] = result.mail_from
            dt['mail_content'] = result.mail_content
            dt['create_date'] = result.create_date.strftime("%Y-%m-%d %H:%M:%S")
            dt['img_path'] = result.img_path
            dt['mark'] = mark
            msg = dt
            # 对应的用户消息数量减一
            if result.is_read == 0:
                obj = systemmessage.objects.get(username=request.user.username)
                obj.unread_mess_mun -= 1
                obj.save()
            # 将对应id的消息置为已读(0表示未读，1表示已读）
            result.is_read = 1
            result.save()


            return Response(msg)
        elif mark == '1':
            result1 = maillogs.objects.get(id=id)
            result2 = monitor.objects.get(id=result1.monitordetail_id)
            dt = {}
            dt['mail_title'] = result1.mail_title
            dt['mail_from'] = result1.mail_from
            dt['mail_content'] = result1.mail_content
            dt['create_date'] = result1.create_date.strftime("%Y-%m-%d %H:%M:%S")
            dt['advice_content'] = result2.proposal_content
            dt['msg_type'] = result2.warning_type
            dt['tagname'] = result2.warning_name
            dt['mark'] = mark
            msg = dt
        if result1.is_read == '0':
            obj = systemmessage.objects.get(username=request.user.username)
            obj.unread_mess_mun -= 1
            obj.save()
        result1.is_read = 1
        result1.save()
    return Response(msg)


@api_view(http_method_names=['GET'])
def msg_delete(request):
    """ 删除对应邮件 """
    if 'id' in request.GET:
        id = request.GET['id']
    maillogs.objects.get(id=id).delete()
    result = {}
    result['code'] = 0
    return Response(result)

@api_view(http_method_names=['GET'])
def get_all_msg_info(request):
    """ 获取所有消息列表 """

    if 'page' in request.GET:
        page = request.GET['page']
        offset = (int(page) - 1) * LIMIT
    else:
        offset = (1) * LIMIT
    if 'msgType' in request.GET:
        msg_type = request.GET['msgType']
    else:
        msg_type = ''
    ordinary_msg = []
    warning_msg = []
    error_msg = []
    if msg_type == 'ordinary' or msg_type == '':
        # 获取普通消息
        where = 'mail_to=' + "'" + request.user.email + "'"
        # 条数
        sql = 'select count(*) from connect_maillogs where ' + where
        # msg_number = utils.getResultBySql(sql)
        # fetchall_ordinary_msg = utils.getResultBySql(
        #     'select id,mail_title,mail_from,mail_content,create_date,img_path,is_read from connect_maillogs where' + ' ' +
        #     where
        # )
        fetchall_ordinary_msg = []
        executeSession = utils.SqlUtils()
        try:
            msg_numberQuerySet = executeSession.getArrResult(sql)
            fetchall_ordinary_msg = executeSession.getArrResult('select id,mail_title,mail_from,mail_content,create_date,img_path,is_read from connect_maillogs where' + ' ' + where)
        except Exception as error:
            executeSession.rollBack()
            logger.error('---error---file:dashboardviews.py;method:get_all_msg_info;error:%s' % error)
        executeSession.closeConnect()
        msg_number = msg_numberQuerySet[0][0] if msg_numberQuerySet else 0

        print('fetchall_ordinary_msg',fetchall_ordinary_msg)
        for obj in fetchall_ordinary_msg:
            dt = {}
            dt['id'] = obj[0]
            dt['mail_title'] = obj[1]
            dt['mail_from'] = obj[2]
            dt['mail_content'] = obj[3]
            dt['create_date'] = obj[4].strftime("%Y-%m-%d %H:%M:%S")
            dt['img_path'] = obj[5]
            dt['is_read'] = obj[6]
            ordinary_msg.append(dt)
        return Response({'msg_number': msg_number, 'ordinary_msg': ordinary_msg})
    if msg_type == 'waring':
        # 获取系统报警消息
        where = 'mail_to=' + "'" + request.user.email + "'" + "and options = '1'"
        # 获取系统报警消息条数
        # msg_number_warning = utils.getResultBySql('select count(*) from connect_maillogs where ' + where)
        # fetchall_warning_msg = utils.getResultBySql(
        #     'select id,mail_title,mail_from,mail_content,create_date,img_path,is_read from connect_maillogs where' + ' ' +
        #     where + ' order by id desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset)
        # )
        fetchall_warning_msg = []
        executeSession = utils.SqlUtils()
        try:
            msg_number_warningQuerySet = executeSession.getArrResult('select count(*) from connect_maillogs where ' + where)
            fetchall_warning_msg = executeSession.getArrResult(
                'select id,mail_title,mail_from,mail_content,create_date,img_path,is_read from connect_maillogs where' + ' ' +
                where + ' order by id desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset))
        except Exception as error:
            executeSession.rollBack()
            logger.error('---error---file:dashboardviews.py;method:get_all_msg_info;error:%s' % error)
        executeSession.closeConnect()
        msg_number_warning = msg_number_warningQuerySet[0][0] if msg_number_warningQuerySet else 0

        for obj in fetchall_warning_msg:
            dt = {}
            dt['id'] = obj[0]
            dt['mail_title'] = obj[1]
            dt['mail_from'] = obj[2]
            dt['mail_content'] = obj[3]
            dt['create_date'] = obj[4].strftime("%Y-%m-%d %H:%M:%S")
            dt['img_path'] = obj[5]
            dt['is_read'] = obj[6]
            warning_msg.append(dt)
        return Response({'msg_number_warning': msg_number_warning, 'warning_msg': warning_msg})
    if msg_type == 'error':
        # 获取系统报错消息
        where = 'mail_to=' + "'" + request.user.email + "'" + "and options = '2'"
        # 条数
        # msg_number_error = utils.getResultBySql('select count(*) from connect_maillogs where ' + where)
        # fetchall_error_msg = utils.getResultBySql(
        #     'select id,mail_title,mail_from,mail_content,create_date,img_path,is_read from connect_maillogs where' + ' ' +
        #     where + ' order by id desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset)
        # )
        fetchall_error_msg = []
        executeSession = utils.SqlUtils()
        try:
            msg_number_warningQuerySet = executeSession.getArrResult('select count(*) from connect_maillogs where ' + where)
            fetchall_error_msg = executeSession.getArrResult(
                'select id,mail_title,mail_from,mail_content,create_date,img_path,is_read from connect_maillogs where' + ' ' +
                where + ' order by id desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset))
        except Exception as error:
            executeSession.rollBack()
            logger.error('---error---file:dashboardviews.py;method:get_all_msg_info;error:%s' % error)
        executeSession.closeConnect()
        msg_number_error = msg_number_warningQuerySet[0][0] if msg_number_warningQuerySet else 0

        for obj in fetchall_error_msg:
            dt = {}
            dt['id'] = obj[0]
            dt['mail_title'] = obj[1]
            dt['mail_from'] = obj[2]
            dt['mail_content'] = obj[3]
            dt['create_date'] = obj[4].strftime("%Y-%m-%d %H:%M:%S")
            dt['img_path'] = obj[5]
            dt['is_read'] = obj[6]
            error_msg.append(dt)
        return Response({'error_msg': error_msg, 'msg_number_error': msg_number_error})

@api_view(http_method_names=['GET','POST'])
def set_read(request):
    """ 将消息一键设置为已读 """
    if 'num' in request.GET:
        num = request.GET['num']
    # 0代表普通消息，1代表报警消息，2代表错误消息
    if num == '0':
        maillogs.objects.filter(mail_to=request.user.email, options='1').update(is_read=1)
        # 消息计数表清零
        obj = systemmessage.objects.get(username=request.user.username)
        obj.unread_mess_mun = 0
        obj.save()
    elif num == '1':
        maillogs.objects.filter(mail_to=request.user.email, options='1').update(is_read=1)
    elif num == '2':
        maillogs.objects.filter(mail_to=request.user.email, options='2').update(is_read=1)
    result = {}
    result['code'] = 0
    return Response(result)


@api_view(http_method_names=['POST'])
def execute_rules(request, id):
    group_left_str = ''  # 分组字段（左）
    select_left_str = ''  # sql select 后的查询字段
    calculation_left_str = ''  # 计算字段（左）
    filter_left_str = ''  # 过滤条件（左）
    sql_str_left = ''  # sql拼接（左）
    use_olap_left = []  # 所使用到的olap
    use_table_left = ''  # 所使用到的表

    ordinary_value = ''  # 右边为普通值
    group_right_str = ''  # 分组字段（右）
    select_right_str = ''  # sql select 后的查询字段
    calculation_right_str = ''  # 计算字段（右）
    filter_right_str = ''  # 过滤条件（右）
    sql_str_right = ''  # sql拼接（右）
    use_olap_right = []  # 所使用到的olap
    use_table_right = ''  # 所使用到的表
    contrast_mode = ''  # 对比方式
    df_use_on = []
    df_use_date = ''  # 用于同比环比日期字段

    # 通过id查询业务规则
    monitorobj = monitor.objects.get(id=id)
    # 获取预警配置信息
    waring_info = monitorobj.waring_calculation_info
    if not waring_info:
        waring_info = {}
    if type(waring_info) == type('stringtype'):
        waring_info = json.loads(waring_info.replace("'", '"'))

    # 左
    # 处理分组字段 默认汇总
    quick = waring_info['quick']
    for group_left in waring_info['group_column_left']:
        if quick in ['y_yp', 'm_yp', 'mp']:  # 同环比日期字段处理
            select_left_str = ' date("' + group_left['col'] + '") as "' + group_left['col'] + '",'
            group_left_str = ' date("' + group_left['col'] + '")'
            df_use_date = group_left['col']
        else:
            group_left_str = group_left_str + '"' + group_left['col'] + '"' + ','
            select_left_str = group_left_str
        df_use_on.append(group_left['col'])

        use_olap_left.append(group_left['olapid'])
    if group_left_str != '':
        group_left_str = ' group by ' + group_left_str
    else:  # 无分组字段
        group_left_str = ''
        select_left_str = ''
    # 处理计算字段 添加 case when 处理null相加问题
    for calculation_left in waring_info['calculation_column_left']:
        if 'col' in calculation_left:  # 字段
            use_olap_left.append(calculation_left['olapid'])
            if quick == 'sum' or quick == '':
                calculation_left_str = calculation_left_str + ' sum( case when "' + calculation_left[
                    'col'] + '" IS NULL THEN 0 ELSE ' + '"' + calculation_left['col'] + '" end )'
            elif quick == 'min':
                calculation_left_str = calculation_left_str + ' min( case when "' + calculation_left[
                    'col'] + '" IS NULL THEN 0 ELSE ' + '"' + calculation_left['col'] + '" end )'
            elif quick == 'max':
                calculation_left_str = calculation_left_str + ' max( case when "' + calculation_left[
                    'col'] + '" IS NULL THEN 0 ELSE ' + '"' + calculation_left['col'] + '" end )'
            else:  # 同环比
                calculation_left_str = ' sum("' + calculation_left['col'] + '")'
        else:  # 运算符
            calculation_left_str = calculation_left_str + ' ' + calculation_left + ' '
    # 处理过滤条件
    for filter_left in waring_info['filter_condition_left']:
        if 'col' in filter_left:
            use_olap_left.append(filter_left['olapid'])
            filter_left_str = filter_left_str + ' where ' + '"' + filter_left['col'] + '"'
        else:
            filter_left_str = filter_left_str + filter_left
    for olapid in set(use_olap_left):
        obj = olap.objects.get(id=olapid)
        use_table_left = use_table_left + obj.table + ','
    use_table_left = use_table_left[:-1]
    if group_left_str.endswith(','):
        group_left_str = group_left_str[:-1]
    sql_str_left = 'select ' + select_left_str + ' ' + calculation_left_str + ' as calculation_left from ' + use_table_left + ' ' + filter_left_str + ' ' + group_left_str
    if waring_info['ordinary_value'] == '':
        # 右
        # 处理分组字段 默认汇总
        for group_right in waring_info['group_column_right']:
            use_olap_right.append(group_right['olapid'])
            group_right_str = group_right_str + '"' + group_right['col'] + '"' + ','
        if group_right_str[:-1] != '':
            select_right_str = group_right_str
            group_right_str = ' group by ' + group_right_str[:-1]  # 去除逗号
        else:  # 无分组字段
            group_right_str = ''
            select_right_str = ''
        # 处理计算字段 添加 case when 处理null相加问题
        for calculation_right in waring_info['calculation_column_right']:
            if 'col' in calculation_right:  # 字段
                use_olap_right.append(calculation_right['olapid'])
                if quick == 'sum' or quick == '':
                    calculation_right_str = calculation_right_str + ' sum( case when "' + calculation_right[
                        'col'] + '" IS NULL THEN 0 ELSE ' + '"' + calculation_right['col'] + '"end )'
                elif quick == 'min':
                    calculation_right_str = calculation_right_str + ' min( case when "' + calculation_right[
                        'col'] + '" IS NULL THEN 0 ELSE ' + '"' + calculation_right['col'] + '"end )'
                elif quick == 'max':
                    calculation_right_str = calculation_right_str + ' max( case when "' + calculation_right[
                        'col'] + '" IS NULL THEN 0 ELSE ' + '"' + calculation_right['col'] + '"end )'
            else:  # 运算符
                calculation_right_str = calculation_right_str + ' ' + calculation_right + ' '
        # 处理过滤条件
        for filter_right in waring_info['filter_condition_right']:
            if 'col' in filter_right:
                use_olap_right.append(filter_right['olapid'])
                filter_right_str = filter_right_str + ' where ' + '"' + filter_right['col'] + '"'
            else:
                filter_right_str = filter_right_str + filter_right
        for olapid in set(use_olap_right):
            obj = olap.objects.get(id=olapid)
            use_table_right = use_table_right + obj.table + ','
        use_table_right = use_table_right[:-1]
        if group_right_str.endswith(','):
            group_right_str = group_right_str[:-1]
        sql_str_right = 'select ' + select_right_str + ' ' + calculation_right_str + ' as calculation_right from ' + use_table_right + ' ' + filter_right_str + ' ' + group_right_str
    else:
        ordinary_value = waring_info['ordinary_value']
    # 执行sql
    fetchall_left = []
    fetchall_right = []
    executeSession = utils.SqlUtils(DBTypeCode.EXTENSION_DB.value)
    try:
        fetchall_left = executeSession.getDictResult(sql_str_left)
        if sql_str_right != '':
            fetchall_right = executeSession.getDictResult(sql_str_right)
    except Exception as error:
        executeSession.rollBack()  # 如果是执行update的sql就需要rollback
        logger.error('---error---file:olapview.py;method:monitor_algorithm;error:%s' % error)  # 需要修改这里的打印信息
    executeSession.closeConnect()
    # 创建pands
    df_left = pd.DataFrame(fetchall_left)
    if fetchall_right:
        df_right = pd.DataFrame(fetchall_right)
    # 开始对比
    contrast_mode = waring_info['contrast_mode']
    # 处理同环比
    yp_mp_df = pd.DataFrame()
    if quick in ['y_yp', 'm_yp', 'mp']:
        # 转换为日期
        df_left = df_left.dropna()  # 去除日期为空的数据
        df_left[df_use_date] = pd.to_datetime(df_left[df_use_date], errors='ignore')
        df_left['year'] = df_left[df_use_date].dt.year  # 年份
        df_left['month'] = df_left[df_use_date].dt.month  # 月份

        # 优先获取过滤条件中的年份信息，如果没有过滤条件，则直接选择df中年份最大的日期，也就是最接近当前的日期
        max_year = df_left['year'].max()  # 获取最大的年份

        df_that_year = df_left[df_left['year'] == max_year]  # 当前年的数据
        df_last_year = df_left[df_left['year'] == (int(max_year) - 1)]  # 去年的数据
        if quick == 'y_yp':  # 年同比
            df_that_year = df_that_year.groupby(['year']).agg(sum).reset_index()
            df_last_year = df_last_year.groupby(['year']).agg(sum).reset_index()
            yp_mp_df['calculation_left_x'] = df_that_year['calculation_left']
            yp_mp_df['calculation_left_y'] = df_last_year['calculation_left']
            yp_mp_df['calculation_left'] = ((yp_mp_df['calculation_left_x'] - yp_mp_df['calculation_left_y']) /
                                            yp_mp_df['calculation_left_y']) * 100
        if quick == 'm_yp':  # 月同比
            current_month = int(datetime.datetime.now().strftime('%m'))
            current_month = 11  # 获取当前月份
            df_that_year = df_that_year.groupby(['year', 'month']).agg(sum).reset_index()
            df_last_year = df_last_year.groupby(['year', 'month']).agg(sum).reset_index()
            df_that_month = df_that_year[df_that_year['month'] == current_month]
            df_last_month = df_last_year[df_that_year['month'] == current_month]
            yp_mp_df['calculation_left_x'] = df_that_month['calculation_left']
            yp_mp_df['calculation_left_y'] = df_last_month['calculation_left']
            yp_mp_df['calculation_left'] = ((yp_mp_df['calculation_left_x'] - yp_mp_df['calculation_left_y']) /
                                            yp_mp_df['calculation_left_y']) * 100
        if quick == 'mp':  # 环比
            current_month = int(datetime.datetime.now().strftime('%m'))
            current_month = 11  # 获取当前月份
            df_that_year = df_that_year.groupby(['year', 'month']).agg(sum).reset_index()
            df_that_month = df_that_year[df_that_year['month'] == current_month]
            df_last_month = df_that_year[df_that_year['month'] == current_month - 1]
            # yp_mp_df['calculation_left_y'] = df_last_month['calculation_left']
            # yp_mp_df['calculation_left_x'] = df_that_month['calculation_left']
            yp_mp_df = df_that_year[df_that_year['month'] >= (current_month - 1)]
            yp_mp_df['calculation_left_y'] = yp_mp_df['calculation_left'].shift()
            yp_mp_df = yp_mp_df.dropna()
            yp_mp_df['calculation_left'] = ((yp_mp_df['calculation_left'] - yp_mp_df['calculation_left_y']) / yp_mp_df[
                'calculation_left_y']) * 100
    else:
        if ordinary_value != '':  # 右边为普通值
            if contrast_mode == 'in':
                df_left = df_left[df_left['calculation_left'].isin(ordinary_value)]
            elif contrast_mode == 'not in':
                df_left = df_left[-df_left['calculation_left'].isin(ordinary_value)]
            elif contrast_mode == '=':
                ordinary_value = int(ordinary_value)
                df_left = df_left[df_left['calculation_left'] == ordinary_value]
            elif contrast_mode == '<':
                ordinary_value = int(ordinary_value)
                df_left = df_left[df_left['calculation_left'] < ordinary_value]
            elif contrast_mode == '<=':
                ordinary_value = int(ordinary_value)
                df_left = df_left[df_left['calculation_left'] <= ordinary_value]
            elif contrast_mode == '>':
                ordinary_value = int(ordinary_value)
                df_left = df_left[df_left['calculation_left'] > ordinary_value]
            elif contrast_mode == '>=':
                ordinary_value = int(ordinary_value)
                df_left = df_left[df_left['calculation_left'] >= ordinary_value]
        else:  # 右边为指标
            df_left = pd.merge(df_left, df_right, on=df_use_on, how='left')
            if contrast_mode == '=':
                df_left = df_left[df_left['calculation_left'] == df_left['calculation_right']]
            elif contrast_mode == '<':
                df_left = df_left[df_left['calculation_left'] < df_left['calculation_right']]
            elif contrast_mode == '<=':
                df_left = df_left[df_left['calculation_left'] <= df_left['calculation_right']]
            elif contrast_mode == '>':
                df_left = df_left[df_left['calculation_left'] > df_left['calculation_right']]
            elif contrast_mode == '>=':
                df_left = df_left[df_left['calculation_left'] >= df_left['calculation_right']]
    if df_left.empty and yp_mp_df.empty:
        pass
    else:  # 如果有值发送消息入库
        mail_content = r'</table></div></br></br>说明/意见:</br>' + monitorobj.proposal_content
        mail_title = '数据眼系统业务规则触发提醒邮件.....'
        # 取到收件人
        mail_to = ''
        if monitorobj.addressee is not None and monitorobj.addressee != '':
            send_user_str = monitorobj.addressee
            select_sql = 'SELECT email FROM account_sys_userextension WHERE username IN (' + send_user_str + ')'
            # emails = utils.getResultBySql(select_sql)
            emails = utils.SqlUtils().getArrResultWrapper(select_sql, logger, 'dashboardviews.py', 'execute_rules')
            if len(emails) > 0:
                temp_email = []
                for mail_row in emails:
                    temp_email.append(mail_row[0])
                mail_to = ','.join(temp_email)
        if mail_to != '':
            try:
                sendTextMail(mail_title, mail_content, mail_to,monitordetail_id=id )
            except:
                pass

        # channels向客户端组发送消息(消息可为任意字符串，只起触发效果)触发ajax

    Group("chat").send({"text": "1"})
    result = {}
    result['code'] = 0
    if df_left.empty:
        result['data'] = df_left
    elif yp_mp_df.empty:
        result['data'] = yp_mp_df
    return Response(result)


@api_view(http_method_names=['POST'])
def execute_rules1(request, state=0):
    print('aaaaaaaaaaaaaaa')
    """ 业务规则系统报警 根据所配置的业务规则 使用拼接sql拼接的方法得到超出条件的数据"""
    if 'olap_data_id' in request.data:
        olapid = request.data['olap_data_id']
    else:
        # 通过保存业务规则方法直接调用

        monitor_form = request.data['monitor']
        if len(str(monitor_form['olapid'])) > 0:
            olapid = str(monitor_form['olapid'])

    info = olap.objects.get(id=olapid)
    olap_name = info.name
    olap_type = info.olaptype
    table = info.table
    rules_data = []
    # 使用最新的版本号
    # try:
    #     version = utils.getResultBySql('select max(version)  from ' + ' ' + table, utils.DATAXEXTENSION_DB_CHAR)[0][0]
    # except Exception as e:
    #     print(e)
    try:
        # 取到所有的业务规则
        rules = monitor.objects.filter(olapid=olapid)
        for rule in rules:
            rule_id = rule.id
            indicate_flag = 0
            mail_content = r'''
                            <style> 
                            .table-mail table{border-right:1px solid #000000;border-bottom:1px solid #000000} 
                            .table-mail table td{border-left:1px solid #000000;border-top:1px solid #000000;padding-left:3px} 
                            </style> ''' \
                           + '''<div class="table-mail"><table><tr><td>''' + olap_type \
                           + '''名称</td><td>业务规则</td><td>触发条件名称</td><td>触发条件详情</td><td>触发条数</td><td>触发时间</td></tr>'''

            # 取到业务规则下的触发方式
            rule_details = monitorDetail.objects.filter(monitorid=rule.id)
            # if len(rule_details) == 0:
            #     return ''
            for rule_detail in rule_details:
                if rule_detail.issend == 'y':
                    sql_str = ''
                    currConditionStr = rule_detail.condition
                    currConditionStr = currConditionStr.replace("'", "\"")
                    currConditionStr = currConditionStr.replace("True", "\"y\"")
                    currConditionStr = currConditionStr.replace("False", "\"n\"")
                    currConditionStr = currConditionStr.replace("None", "\"\"")
                    currConditionStr = currConditionStr.replace("\\xa0", " ")
                    ruleList = json.loads(currConditionStr)
                    #sql_str = spliceSql(rule_detail, ruleList, table, version)
                    # 查询出符合触发方式的数据条数
                    #cnt = utils.getResultBySql(sql_str, utils.DATAXEXTENSION_DB_CHAR)[0][0]
                    cnt = 1
                    if int(cnt) > 0 and state != 1:
                        # 拼接需要记录的数据并存入数据库
                        # olapid, olap名称, olap类型, 规则id, 规则名称, 触发条件id，触发条件名称，颜色，sql，条数，版本，当前时间
                        data_str = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                        json_obj = {}
                        monitordetail_id = rule_detail.id
                        json_obj['olap_id'] = olapid
                        json_obj['olap_name'] = olap_name
                        json_obj['olap_type'] = olap_type
                        json_obj['monitor_id'] = rule.id
                        json_obj['monitor_name'] = rule.title
                        json_obj['monitor_detail_id'] = rule_detail.id
                        json_obj['monitor_detail_name'] = rule_detail.tagname
                        json_obj['monitor_detail_color'] = rule_detail.color
                        json_obj['advice_content'] = rule_detail.advice_content
                        json_obj['cnt'] = cnt
                        #json_obj['version'] = version
                        json_obj['date_str'] = data_str
                        begin = sql_str.find('FROM')
                        #objs = utils.getResultBySql('SELECT * ' + sql_str[begin:], utils.DATAXEXTENSION_DB_CHAR)
                        indicate_content = []
                        rules_data = []
                        rules_data.append(json_obj)

                        # for obj in objs:
                        #     rules = {}
                        #     rules['rules_data'] = obj
                        #     indicate_content.append(rules)
                        #     rules_data.append(rules)

                        rulelogs.objects.create(olap_id=olapid, monitor_id=rule.id, monitordetail_id=rule_detail.id
                                                , describe=json.dumps(json_obj), indicate_cnt=int(cnt),
                                                indicate_content=indicate_content)
                        mail_content = mail_content + r'<tr><td>' + olap_name + '</td><td>' + rule.title + '</td><td>' \
                                       + rule_detail.tagname + '</td><td>' + rule_detail.show_condition_str \
                                       + '</td><td>' + str(cnt) + '</td><td>' + data_str + '</td></tr>'
                        indicate_flag = 1
            # indicate_flag = 0
            if indicate_flag == 1 and state != 1:
                # 如果用户配置了邮件模板使用已配置的模板
                emailobj = emailconf.objects.filter()
                # 获取email配置信息。系统分享类模板所使用的模板id
                #templateid = json.loads(emailobj[0].options)['emailrulestemplate']
                #templateobj = msgalltemplateconf.objects.filter(id=templateid)
                # if templateobj:
                #     mail_title = templateobj[0].templatetitle
                #     mail_content = templateobj[0].templatecontent + mail_content + \
                #                    r'</table></div></br></br>说明/意见:</br>' + rule_detail.advice_content
                # else:
                mail_content = mail_content + r'</table></div></br></br>说明/意见:</br>' + rule_detail.advice_content
                mail_title = '数据眼系统业务规则触发提醒邮件.....'
                # 取到收件人
                mail_to = ''
                if rule.receive_user is not None and rule.receive_user != '':
                    send_user_str = rule.receive_user
                    select_sql = 'SELECT email FROM account_sys_userextension WHERE username IN (' + send_user_str + ')'
                    # emails = utils.getResultBySql(select_sql)
                    emails = utils.SqlUtils().getArrResultWrapper(select_sql,logger, 'dashboardviews.py', 'execute_rules')
                    if len(emails) > 0:
                        temp_email = []
                        for mail_row in emails:
                            temp_email.append(mail_row[0])
                        mail_to = ','.join(temp_email)
                # 取到抄送人
                acc = ''
                if rule.cc_user is not None and rule.cc_user != '':
                    cc_user_str = rule.cc_user
                    select_sql = 'SELECT email FROM account_sys_userextension WHERE username IN (' + cc_user_str + ')'
                    # emails = utils.getResultBySql(select_sql)
                    emails = utils.SqlUtils().getArrResultWrapper(select_sql, logger,'dashboardviews.py', 'execute_rules')
                    if len(emails) > 0:
                        temp_email = []
                        for mail_row in emails:
                            temp_email.append(mail_row[0])
                            acc = ','.join(temp_email)
        if mail_to != '':
            sendTextMail(mail_title, mail_content, mail_to, acc, rule_id, monitordetail_id)

    except Exception as e:
        raise e
        print(e)
        print(sys.exc_info())
    # channels向客户端组发送消息(消息可为任意字符串，只起触发效果)触发ajax
    if state != 1:
        Group("chat").send({"text": "1"})
    result = {}
    result['code'] = 0
    result['data'] = rules_data
    return Response(result)


# 拼接sql
def spliceSql(rule_detail, conditions, tableName, version):
    sql = ''
    try:
        # 以其他表中字段为依据进行判断的字段
        exrtaOlapCols = []
        # 在页面维护值得字段
        for condition in conditions:
            exrtaOlapCols.append(condition)
        sql = "SELECT COUNT(*) FROM" + " "
        # 拼接用到的所有表
        sql_table = ''
        # 拼接所有的olap版本
        versionstr = ''
        # 拼接所有的关联关系
        sql_connection = ''
        condition_str = ''
        # where条件
        where_one = ''
        where_two = ''
        for exrtaOlapCol in exrtaOlapCols:
            joinRows = exrtaOlapCol['joinRow']
            for joinRow in joinRows:
                if (joinRow['col1'] == '' and joinRow['aimCol1'] == '') and (
                        joinRow['col2'] == '' and joinRow['aimCol2'] == ''):
                    sql_connection = ''
                elif joinRow['col1'] == '' and joinRow['aimCol1'] == '':
                    sql_connection = joinRow['col2'] + "=" + joinRow['aimCol2']
                elif joinRow['col2'] == '' and joinRow['aimCol2'] == '':
                    sql_connection = joinRow['col1'] + "=" + joinRow['aimCol1']
                else:
                    sql_connection = sql_connection + joinRow['col1'] + "=" + joinRow['aimCol1'] + " " + "and" + \
                                     " " + joinRow['col2'] + "=" + joinRow['aimCol2']
            condition_str = rule_detail.condition_str
            # 业务规则中左右都有where条件
            if condition_str.find('&&&') != -1 and condition_str.find('%%%') != -1:
                temp = condition_str[condition_str.find('&&&') + 3:]
                where_one = temp[:temp.find('&&&')]
                temp = condition_str[condition_str.find('%%%') + 3:]
                where_two = temp[:temp.find('%%%')]
                condition_str = condition_str.replace('&&&' + where_one + '&&&', '').replace('%%%' + where_two + '%%%',
                                                                                             '')
                where_one = ' and ' + where_one
                where_two = ' and ' + where_two
                condition_str = char_to_num(condition_str)
            # 业务规则中左边有where条件
            elif condition_str.find('&&&') != -1 and condition_str.find('%%%') == -1:
                temp = condition_str[condition_str.find('&&&') + 3:]
                where_one = temp[:temp.find('&&&')]
                condition_str = condition_str.replace('&&&' + where_one + '&&&', '')
                where_one = ' and ' + where_one
                condition_str = char_to_num(condition_str)
            # 业务规则中右边有where条件
            elif condition_str.find('&&&') == -1 and condition_str.find('%%%') != -1:
                temp = condition_str[condition_str.find('%%%') + 3:]
                where_two = temp[:temp.find('%%%')]
                condition_str = condition_str.rreplace('%%%' + where_two + '%%%', '')
                where_two = ' and ' + where_two
                condition_str = char_to_num(condition_str)
            # 业务规则两边都没有where条件
            else:
                condition_str = char_to_num(condition_str)

        for table in rule_detail.use_tables.split(','):
            # version = utils.getResultBySql('select max(version)  from ' + ' ' + table,utils.DATAXEXTENSION_DB_CHAR)[0][0]
            versionQuerySet = utils.SqlUtils(DBTypeCode.EXTENSION_DB.value).getArrResultWrapper('select max(version)  from ' + ' ' + table,
                                                                  logger,'dashboardviews.py', 'spliceSql')
            version = versionQuerySet[0][0] if versionQuerySet else 0

            versionstr = versionstr + " " + 'and' + ' ' + table + '.' + 'version' + '=' + str(version)
            sql_table = sql_table + table + ","
        if sql_connection == '':
            sql = sql + " " + sql_table[:-1] + " " + "where" + " " + sql_connection + " " + \
                  " " + condition_str + versionstr + where_one + where_two
        else:
            sql = sql + " " + sql_table[:-1] + " " + "where" + " " + sql_connection + " " + "and" + \
                  " " + condition_str + versionstr + where_one + where_two
    except Exception as e:
        print(e)
    return sql


@api_view(http_method_names=['POST'])
def saveAdditionText(request):
    targetId = request.data['targId']
    # targetType=request.data['type']
    targetText = request.data['text']
    sceneId = request.data['sceneId']
    hideStatus = request.data['hideStatus']
    updateFoledSts = ''
    if 'updateFoledSts' in request.data:
        updateFoledSts = request.data['updateFoledSts']

    rs = {'code': 0}
    try:
        # save to scene
        sceneObj = scenes.objects.get(id=sceneId)
        sceneOption = sceneObj.options

        commentTextObj = {}
        commentTextObj['id'] = targetId
        commentTextObj['html'] = targetText
        commentTextObj['hideStatus'] = hideStatus
        if sceneOption:
            if type(sceneOption) == type('stringtype'):
                sceneOption = json.loads(sceneOption)
            if 'addtionCommonText' not in sceneOption:
                sceneOption['addtionCommonText'] = []
        else:
            sceneOption = {'addtionCommonText': []}

        # 如果传进来的targetText是空就删除,否则如果有id就修改，没有就添加
        indexOfNullObj = 0
        if not targetText or not targetText.strip():
            popCnt = 0
            for commTxObj in sceneOption['addtionCommonText']:
                if commTxObj['id'] == targetId:
                    sceneOption['addtionCommonText'].pop(indexOfNullObj - popCnt)  # 删除
                    popCnt += 1
                    # break;
                else:
                    indexOfNullObj += 1
        else:
            for commTxObj in sceneOption['addtionCommonText']:
                if commTxObj['id'] == targetId:  # 修改
                    if updateFoledSts and updateFoledSts == 'yes':  # 如果只修改展开/关闭状态
                        commentTextObj = sceneOption['addtionCommonText'][indexOfNullObj]
                        commentTextObj['hideStatus'] = hideStatus
                    sceneOption['addtionCommonText'][indexOfNullObj] = commentTextObj
                    break;
                else:
                    indexOfNullObj += 1
            else:
                sceneOption['addtionCommonText'].append(commentTextObj)
        sceneObj.options = json.dumps(sceneOption)
        sceneObj.save()
        rs['code'] = 1
    except Exception as e:
        raise e
        print('==error==file:dashboardview.py;method:saveAdditionText;line:2288')
        rs['msg'] = e.args
    return Response(rs)


# 获取业务规则配置运行后的返回信息
@api_view(http_method_names=['GET'])
def getBussRuleMsg(request):
    rsMsg = {'code': 0}
    try:
        olapId = request.query_params['olapId']
        olapObj = olap.objects.get(id=olapId)
        if not olapObj:
            rsMsg['message'] = '没有找到olap！'
        elif not olapObj.options:
            rsMsg['message'] = '业务规则返回信息为空！'
        else:
            optionObj = olapObj.options
            if type(optionObj) == type('stringtype'):
                optionObj = json.loads(optionObj)
            if 'bussRuleMsg' in optionObj:
                rsMsg['bussRuleMsg'] = optionObj['bussRuleMsg']
                rsMsg['code'] = 1
    except Exception as err:
        print('=====error=====file:dashboardviews.py;method:getBussRuleMsg;line:2314;error Msg=', err.args)
    return Response(rsMsg)


# 统计场景的各个标签
@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def statisticalTags(request):
    results = tools.successMes()
    try:
        tableName = request.query_params['table']
        colName = request.query_params['col']
        sql = '''
            SELECT
                tagtable.tag,
                COUNT (tagtable.tag) AS tagnum
            FROM
                (
                    SELECT
                        regexp_split_to_table(
                            REPLACE ("{0}", '，', ','),
                            ','
                        ) AS tag
                    FROM
                        {1}
                    WHERE
                        {1}.sceneversion >= 2
                ) AS tagtable
            GROUP BY
                tagtable.tag           
        '''
        # queryResults = utils.getResultBySql(sql.format(colName, tableName))
        queryResults = utils.SqlUtils().getArrResultWrapper(sql.format(colName, tableName),
                                                                    logger,'dashboardviews.py', 'statisticalTags')
        tags = []
        for queryResult in queryResults:
            tag = {
                'tag': queryResult[0],
                'num': queryResult[1]
            }
            tags.append(tag)
        results['tags'] = tags
    except Exception as e:
        print('file:dashboardviews; method:statisticalTags')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


# 获取用户上传的图片
@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getUserImgs(request):
    results = tools.successMes()
    try:
        userId = request.user.id
        L = []
        imgPath = './frontend/upload/user_' + str(userId) + '/'
        L = getCustomizeImgs(imgPath)
        results['imgs'] = L
    except Exception as e:
        print('file:dashboardviews; method:getUserImgs')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)

# 获取用户上传的组件背景图片
@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def getUserChartsImgs(request):
    results = tools.successMes()
    try:
        userId = request.user.id
        L = []
        imgPath = './frontend/upload/user_' + str(userId) + '/chartsbgimgs/'
        L = getCustomizeImgs(imgPath)
        results['imgs'] = L
    except Exception as e:
        print('file:dashboardviews; method:getUserChartsImgs')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


# 删除用户上传的图片
@api_view(http_method_names=['GET'])
@permission_classes((permissions.AllowAny,))
def delUserImg(request):
    results = tools.successMes()
    try:
        imgPath = request.GET['imgPath']
        # 删除
        delPath = '.' + imgPath
        delCustomizeImgs(delPath)
    except Exception as e:
        print('file:dashboardviews; method:delUserImgs')
        print(e)
        return Response(tools.errorMes(e.args))
    return Response(results)


def delCustomizeImgs(imgPath):
    if os.path.exists(imgPath):
        os.remove(imgPath)


def getCustomizeImgs(imgPath):
    paths = []
    if not os.path.exists(imgPath):
        os.makedirs(imgPath)
    for root, dirs, files in os.walk(imgPath):
        for file in files:
            if ('.jpeg' in os.path.splitext(file)[1] or
                '.jpg' in os.path.splitext(file)[1] or
                '.png' in os.path.splitext(file)[1] or
                '.bmp' in os.path.splitext(file)[1] or
                '.gif' in os.path.splitext(file)[1]) and 'default.' not in file:
                paths.append(os.path.join(root, file))
    return paths


def char_to_num(condition_str):
    """ 解决sql字符串类型无法进行计算的问题 """
    condition_str_temp = condition_str.replace('(', '').replace(')', '').replace('+', ' ').replace('-', ' ') \
        .replace('*', ' ').replace('/', ' ').replace('<=', ' ').replace('<', ' ').replace('>=', ' ') \
        .replace('>', ' ').replace('=', ' ').split(' ')
    for temp1 in condition_str_temp:
        replace = 'cast(' + temp1 + ' as NUMERIC)'
        condition_str = condition_str.replace(temp1, replace)
    return condition_str


@api_view(http_method_names=['POST'])
def save_msg_config(request):
    """ 保存消息配置信息 ,options字段存储各个消息类型所选择的模版内容。"""
    if 'id' in request.GET:
        id = request.GET['id']
    # 消息类型为邮件
    if request.data['email']:
        tempobj = emailconf.objects.filter()
        obj = request.data['email']
        options = {}
        options['emailtemplateconf'] = obj['emailtemplateconf']
        options['emailsharetemplate'] = obj['emailsharetemplate']
        options['emailrulestemplate'] = obj['emailrulestemplate']

        if tempobj:
            emailconf.objects.filter().update(mailaddress=obj['mailaddress'], mailpassword=obj['mailpassword']
                                              , smptaddress=obj['smptaddress'], smptport=obj['smptport']
                                              , options=json.dumps(options))
        else:
            emailconf.objects.create(mailaddress=obj['mailaddress'], mailpassword=obj['mailpassword']
                                     , smptaddress=obj['smptaddress'], smptport=obj['smptport'],
                                     options=json.dumps(options))
    # 消息类型为短信
    if request.data['sms']:
        tempobj = smsconf.objects.filter()
        obj = request.data['sms']
        options = {}
        options['smstemplateconf'] = obj['smstemplateconf']
        options['smssharetemplate'] = obj['smssharetemplate']
        options['smsrulestemplate'] = obj['smsrulestemplate']
        if tempobj:
            smsconf.objects.filter().update(appid=obj['sppid'], appkey=obj['appkey'], templateid=obj['templateid'],
                                            template=obj['template'], options=json.dumps(options))
        else:
            smsconf.objects.create(appid=obj['appid'], appkey=obj['appkey'], templateid=obj['templateid'],
                                   template=obj['template'], options=json.dumps(options))
    # 系统配置信息
    if request.data['sys']:
        obj = request.data['sys']
        objs = syserrorconf.objects.filter()
        if objs:
            syserrorconf.objects.update(is_use=obj['msg'], is_email=obj['email'], is_sms=obj['sms'],
                                        is_wechat=obj['wechat'], candidate_email=obj['email_user'],
                                        candidate_sms=obj['sms_user'],
                                        email_template=obj['emailtemplate'], sms_template=obj['smstemplate'],
                                        wechat_template=obj['wechattemplate'])
        else:
            syserrorconf.objects.create(is_use=obj['msg'], is_email=obj['email'], is_sms=obj['sms'],
                                        is_wechat=obj['wechat']
                                        , candidate_email=obj['email_user'], candidate_sms=obj['sms_user'],
                                        email_template=obj['emailtemplate'], sms_template=obj['smstemplate'],
                                        wechat_template=obj['wechattemplate'])
    # 消息类型为微信
    if request.data['wx']:
        obj = request.data['wx']
        options = {}
        options['wxtemplateconf'] = obj['wxtemplateconf']
        options['wxsharetemplate'] = obj['wxsharetemplate']
        options['wxrulestemplate'] = obj['wxrulestemplate']

        tempobj = wechatconf.objects.filter()
        if tempobj:
            wechatconf.objects.filter().update(userid=obj['name'], options=json.dumps(options))
        else:
            wechatconf.objects.create(userid=obj['name'], options=json.dumps(options))
    if request.data['template']:
        obj = request.data['template']
        tempobj = msgalltemplateconf.objects.filter(id=id)
        if tempobj:
            msgalltemplateconf.objects.filter(id=id).update(templatename=obj['templatename'],
                                                            templatetitle=obj['templatetitle'],
                                                            templatecontent=obj['templatecontent'])
        else:
            msgalltemplateconf.objects.create(templatename=obj['templatename'], templatetitle=obj['templatetitle'],
                                              templatecontent=obj['templatecontent'])
    result = {}
    result['code'] = 0
    return Response(result)


@api_view(http_method_names=['GET'])
def get_msg_conf(request):
    """ 此方法主要供消息模板配置页面使用。 获取所有消息配置列表 """
    templateconfs = []
    where = ''
    if 'page' in request.GET:
        page = request.GET['page']
        offset = (int(page) - 1) * LIMIT
    # allcnt = utils.getResultBySql('select count(*) from dashboard_msgalltemplateconf' + where)
    # templateobjs = utils.getResultBySql(
    #     'select templatename,templatename,templatename,status,id from dashboard_msgalltemplateconf' +
    #     where + ' order by status desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset)
    # )
    templateobjs = []
    executeSession = utils.SqlUtils()
    try:
        allcntQuerySet = executeSession.getArrResult('select count(*) from dashboard_msgalltemplateconf' + where)
        templateobjs = executeSession.getArrResult(
            'select templatename,templatename,templatename,status,id from dashboard_msgalltemplateconf' +
             where + ' order by status desc LIMIT ' + str(LIMIT) + ' offset ' + str(offset))
    except Exception as error:
        executeSession.rollBack()
        logger.error('---error---file:dashboardviews.py;method:get_msg_conf;error:%s' % error)
    executeSession.closeConnect()
    allcnt = allcntQuerySet[0][0] if allcntQuerySet else 0


    for obj in templateobjs:
        dt = {}
        dt['templatename'] = obj[0]
        dt['templatetitle'] = obj[1]
        dt['templatecontent'] = obj[2]
        dt['status'] = obj[3]
        dt['id'] = obj[4]
        templateconfs.append(dt)
    result = {}
    result['code'] = 0
    result['templateconfs'] = templateconfs
    result['total'] = allcnt
    return Response(result)


@api_view(http_method_names=['GET'])
def echo_msg_conf(request):
    """ 回显信息配置数据 用于消息参数配置页面"""
    if 'id' in request.GET:
        id = request.GET['id']
    if 'type' in request.GET:
        type = request.GET['type']
    dt = {}
    # email配置
    if type == 'email':
        obj = emailconf.objects.filter()
        options = json.loads(obj[0].options)
        if obj:
            dt['mailaddress'] = obj[0].mailaddress
            dt['mailpassword'] = obj[0].mailpassword
            dt['smptaddress'] = obj[0].smptaddress
            dt['smptport'] = obj[0].smptport
            dt['emailtemplateconf'] = options['emailtemplateconf']
            dt['emailsharetemplate'] = options['emailsharetemplate']
            dt['emailrulestemplate'] = options['emailrulestemplate']
    # 短信配置
    elif type == 'sms':
        obj = smsconf.objects.filter()
        options = json.loads(obj[0].options)
        if obj:
            dt['sppid'] = obj[0].appid
            dt['appkey'] = obj[0].appkey
            dt['templateid'] = obj[0].templateid
            dt['template'] = obj[0].template
            dt['smstemplateconf'] = options['smstemplateconf']
            dt['smssharetemplate'] = options['smssharetemplate']
            dt['smsrulestemplate'] = options['smsrulestemplate']
    # 消息模板配置页面中用来回显数据。
    elif type == 'emailtemplate':
        if id:
            obj = msgalltemplateconf.objects.filter(id=id)
            if obj:
                dt['templatename'] = obj[0].templatename
                dt['templatetitle'] = obj[0].templatetitle
                dt['templatecontent'] = obj[0].templatecontent
        else:
            obj = msgalltemplateconf.objects.filter()
            if obj:
                dt['templatename'] = obj[0].templatename
                dt['templatetitle'] = obj[0].templatetitle
                dt['templatecontent'] = obj[0].templatecontent
    # 微信配置
    elif type == 'wx':
        obj = wechatconf.objects.filter()
        options = json.loads(obj[0].options)
        if obj:
            dt['name'] = obj[0].userid
            dt['wxtemplateconf'] = options['wxtemplateconf']
            dt['wxsharetemplate'] = options['wxsharetemplate']
            dt['wxrulestemplate'] = options['wxrulestemplate']
    # 系统错误提醒配置
    elif type == 'sys':
        sysobj = syserrorconf.objects.filter()
        if sysobj:
            dt['msg'] = sysobj[0].is_use
            dt['email'] = sysobj[0].is_email
            dt['sms'] = sysobj[0].is_sms
            dt['wechat'] = sysobj[0].is_wechat
            dt['email_user'] = sysobj[0].candidate_email
            dt['sms_user'] = sysobj[0].candidate_sms
            dt['emailtemplate'] = sysobj[0].email_template
            dt['smstemplate'] = sysobj[0].sms_template
            dt['wechattemplate'] = sysobj[0].wechat_template
    return Response(dt)


@api_view(http_method_names=['GET'])
def del_msg_template(request):
    """ 删除消息配置模板"""
    if 'id' in request.GET:
        id = request.GET['id']
    if 'type' in request.GET:
        type = request.GET['type']
    if type == 'emailtemplate':
        msgalltemplateconf.objects.filter(id=id).delete()
    result = {}
    result['code'] = 0
    return Response(result)


@api_view(http_method_names=['GET'])
def get_login_page_info(request):
    """ 获取登录页配置信息 """
    objs = loginpageconf.objects.all()
    loginobj = []
    for obj in objs:
        dt = {}
        dt['id'] = obj.id
        dt['is_use'] = obj.is_use
        dt['use_type'] = obj.use_type
        dt['showimage'] = obj.showimage
        loginobj.append(dt)
    result = {}
    result['code'] = 0
    return Response(loginobj)


@api_view(http_method_names=['GET'])
def set_login_page(request):
    """ 设置前后台登录页 1表示后台登录页 ，2表示前台登录页 ，3表示前后台登录页"""
    if 'id' in request.GET:
        id = request.GET['id']
    if 'type' in request.GET:
        # 存储前台传递过来的页面设置值。
        type = request.GET['type']
    if id:
        # 处理在设置之前存在一个登录页同时设置前后台的方式
        obj = loginpageconf.objects.filter(use_type=3)
        if obj:
            # 处理同时设置前后台登录页问题
            if type == '3':
                # 如果id为同一个则直接跳过
                if obj[0].id == id:
                    pass
                else:
                    # 所有的状态恢复默认值（未使用状态）
                    loginpageconf.objects.update(use_type=0, is_use=0)
                    # 获取对应id设置对应参数
                    obj = loginpageconf.objects.get(id=id)
                    obj.is_use = 1
                    obj.use_type = type
                    obj.save()
            # 处理非同时设置前后台登录页的问题。
            else:
                # loginpageconf.objects.update(use_type=0, is_use=0)
                # 更改在设置之前存在一个登录页同时设置前后台的对象参数
                if type == '1':
                    obj.update(use_type=2)
                else:
                    obj.update(use_type=1)
                # 设置新登录页的对象参数
                obj = loginpageconf.objects.get(id=id)
                obj.is_use = 1
                obj.use_type = type
                obj.save()
        # 在设置之前不存在一个登录页同时设置前后台的形式。
        else:
            # 在设置之前是否存在已使用的登录页
            objs = loginpageconf.objects.filter(is_use=1)
            if objs:
                for obj in objs:
                    # 如果id相同，且类型不同，此时设置属性为3，代表前后台同一个登录页
                    if obj.id == id and obj.use_type != type:
                        loginpageconf.objects.update(use_type=0, is_use=0)
                        obj.use_type = 3
                        obj.save()
                    else:
                        # 如果类型不同则直接设置新的登录页
                        if obj.use_type != type:
                            loginpageconf.objects.filter(use_type=type).update(use_type=0, is_use=0)
                            obj = loginpageconf.objects.get(id=id)
                            obj.is_use = 1
                            obj.use_type = type
                            obj.save()
                        # 如果类型相同，则先清除之前的设置，在设置新的登录页。
                        else:
                            loginpageconf.objects.update(use_type=0, is_use=0)
                            obj = loginpageconf.objects.get(id=id)
                            obj.is_use = 1
                            obj.use_type = type
                            obj.save()
            # 在设置之前没设置过登录页
            else:
                loginpageconf.objects.update(use_type=0, is_use=0)
                obj = loginpageconf.objects.get(id=id)
                obj.is_use = 1
                obj.use_type = type
                obj.save()
    result = {}
    result['code'] = 0
    return Response(result)


# 上传文件并存入数据库
@api_view(http_method_names=['POST'])
def upload_uncompress_zipfile(request):
    rs = {}
    f = request.FILES['filename']
    filename = f.name
    # 保存文件到/frontend/upload/temp_files
    saveDir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    pgConfFiledir = os.path.join(saveDir, 'frontend', 'loginpackage')
    dirPath = pgConfFiledir + os.sep + filename
    savepath = pgConfFiledir + os.sep
    with open(dirPath, 'wb') as fw:  # 写入到指定目录
        fw.write(f.read())
    # 解压zip文件
    with zipfile.ZipFile(dirPath, 'r') as zips:
        zips.extractall(savepath)
        zips.close()
    os.remove(dirPath)
    # 信息入库
    showimage = '/frontend/loginpackage/' + filename[:-4] + '/image/login.jpg'
    homepage = 'loginpackage/' + filename[:-4] + '/login.html'
    loginpageconf.objects.create(showimage=showimage, homepage=homepage)
    try:
        rs,
    except Exception as e:
        rs['code'] = 0
        rs['msg'] = e.args
        raise e
    return Response(rs)


# 下载文件
@api_view(http_method_names=['GET'])
def download_file(request):
    saveDir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    pgConfFiledir = os.path.join(saveDir, 'frontend', 'loginpackage')
    dirPath = pgConfFiledir + os.sep + 'logintemplate.zip'
    file_name = '/frontend/loginpackage/logintemplate.zip'

    def file_iterator(file_name, chunk_size=512):
        with open(dirPath, 'rb') as f:
            while True:
                c = f.read(chunk_size)
                if c:
                    yield c
                else:
                    break

    the_file_name = "logintemplate.zip"
    response = StreamingHttpResponse(file_iterator(the_file_name))
    print(response)
    return response


@api_view(http_method_names=['GET'])
def getScenesLogList(request):
    '''
    用户访问场景记录
    '''
    resultObj = tools.successMes()
    try:
        parameters = ''
        if 'first_name' in request.query_params:
            parameters += ' AND asu.first_name like \'%'+request.query_params['first_name']+'%\' '
        if 'scenes_name' in request.query_params:
            parameters += ' AND ds."name" like \'%'+request.query_params['scenes_name']+'%\' '
        if 'keywords' in request.query_params:
            parameters += ' AND ds.keywords like \'%'+request.query_params['keywords']+'%\' '
        if 'modifytimestart' in request.query_params:
            parameters += ' AND du.modifytime >= \''+request.query_params['modifytimestart']+' 00:00:00\' '
        if 'modifytimeend' in request.query_params:
            parameters += ' AND du.modifytime <= \''+request.query_params['modifytimeend']+' 23:59:59\' '
        executeSql = '''
            SELECT
                du."id",
                asu.username,
                asu.first_name,
                ds."name" AS scenes_name,
                ds.keywords,
                du.targettype,
                du.visitcount,
                du.modifytime 
            FROM
                dashboard_uservisitlog du
                LEFT JOIN dashboard_scenes ds ON du.targetid = ds."id"
                LEFT JOIN account_sys_userextension asu ON asu."id" = du.userid 
            WHERE
                1 = 1 {0}
            ORDER BY
                du.modifytime DESC
        '''
        executeSql = executeSql.format(parameters)
        excuteSession = utils.SqlUtils()  
        resultList = excuteSession.getDictResult(executeSql)
        resultTotal = excuteSession.getTotal(executeSql)
        resultObj['data']['list'] = resultList
        resultObj['data']['total'] = resultTotal
        excuteSession.closeConnect()
    except Exception as error:
        print('---api,dashboardapi,dashboardviews.py getScenesLogList-error=', error)
        if excuteSession:
            excuteSession.rollBack()
        resultObj = tools.errorMes(error.args)
    return Response(resultObj)

# 用户访问场景记录
@api_view(http_method_names=['GET'])
def getUserVisitLogList(request):
    resultObj = tools.successMes()
    try:
        currUser = request.user
        currUserObj = sys_userextension.objects.get(id=currUser.id)
        executeSql = 'select du.id,asu.username,ds.name,du.targettype,du.visitcount,du.modifytime from dashboard_uservisitlog du ' \
                     'left join dashboard_scenes ds on du.targetid = ds.id left join account_sys_userextension asu ' \
                     'on asu.id = du.userid '
        where = ' where 1=1 '
        print('currUserObj.is_superuser=',currUserObj.is_superuser)
        if not currUserObj.is_superuser:
            where += " and du.userid = '" + str(currUser.id) + "'"
        if 'search' in request.GET:
            where += """ and ds.name like '%%""" + request.GET['search'] + """%%' """
        executeSql += where #把where条件拼接上
        if 'page' in request.GET:
            offset = (int(request.GET['page']) - 1) * LIMIT
            executeSql = executeSql + """ limit """ + str(LIMIT) + """ offset """ + str(offset)
        # userVisitLogList = utils.getResultBySql(executeSql)
        userVisitLogList = utils.SqlUtils().getArrResultWrapper(executeSql,logger,'dashboardviews.py', 'getDataTable')
        userVisitLogListObjs = []
        for visitLog in userVisitLogList :
            tempVisitLogDict = {}
            tempVisitLogDict['id'] = visitLog[0]
            tempVisitLogDict['username'] = visitLog[1]
            tempVisitLogDict['targetname'] = visitLog[2]
            tempVisitLogDict['targettype'] = visitLog[3]
            tempVisitLogDict['visitcount'] = visitLog[4]
            tempVisitLogDict['modifytime'] = visitLog[5].strftime("%Y-%m-%d %H:%M:%S")
            userVisitLogListObjs.append(tempVisitLogDict)
        resultObj['data'] = userVisitLogListObjs
        resultObj['totalCount'] = len(userVisitLogListObjs)
    except Exception as error:
        print('---getUserVisitLogList-error=',error)
        resultObj = tools.errorMes(error.args)
    return Response(resultObj)

#删除用户访问场景记录by id
@api_view(http_method_names=['GET'])
def deleteUserVisitLogById(request):
    resultObj = tools.successMes()
    try:
        uvId = request.GET['id']
        currUserVisitLogObj = uservisitlog.objects.get(id=uvId)
        currUserVisitLogObj.delete()
    except Exception as error:
        print('---deleteUserVisitLogById-error=',error)
        resultObj = tools.errorMes(error.args)
    return Response(resultObj)


# 保存数据字典
@api_view(http_method_names=['GET','POST'])
def save_data_dictionary(request):
    # 获取添加数据
    add_dictionary = request.data['addedNodes']
    #获取更新数据
    update_dictionary = request.data['updatedNodes']
    print('add_dictionary',add_dictionary)
    print('update_dictionary',update_dictionary)

    try:
        # 循环更新数据
        for dictionary in update_dictionary:
            obj = data_dictionary.objects.get(id=str(dictionary['id']))
            obj.dictionary_name=dictionary['dictionary_name']
            obj.order_num=dictionary['order_num']
            obj.status=dictionary['status']
            obj.save()

        # 添加数据循环入库
        for dictionary in add_dictionary:
            # 统一入库
            if dictionary['parent_id'] == '':
                dictionary['parent_id'] = 0
            data_dictionary.objects.create(dictionary_name=dictionary['dictionary_name'],parent_id=dictionary['parent_id'],
                                         order_num=dictionary['order_num'],code=dictionary['code'],description=dictionary['description']
                                         ,status=dictionary['status'])

    except Exception as err:
        error = err
    result = {}
    result['code'] = 0
    return Response(result)


# 删除数据字典
@api_view(http_method_names=['POST'])
def delete_data_dictionary(request):
    # 获取删除字典数组
    del_dictionary = request.data['deleteNode']
    try:
        # 获取parent_id
        parent_id = del_dictionary['code']
        delete_dictionary(parent_id)
    except Exception as err:
        error = err
    result = {}
    result['code'] = 0
    return Response(result)


# 显示数据字典
@api_view(http_method_names=['GET'])
def show_data_dictionary(request):
    # 获取父节点id
    if 'id' in request.GET:
        id = request.GET['id']
    # 取到所有根节点
    parent_dictionary = []
    firt_objs = data_dictionary.objects.filter(parent_id=0)
    try:
        for obj in firt_objs:
            dist = {}
            dist['id'] = obj.id
            dist['dictionary_name'] = obj.dictionary_name
            dist['parent_id'] = obj.parent_id
            dist['order_num'] = obj.order_num
            dist['code'] = obj.code
            dist['status'] = obj.status
            dist['description'] = obj.description
            dist['nodes'] = []
            second_objs = data_dictionary.objects.filter(parent_id=obj.code)
            if second_objs:
                parent_dictionary.append(two_level_dictionary(dist))
            else:
                parent_dictionary.append(dist)
    except Exception as err:
        logger.error(err)
        print(err)

    result = {}
    result['code'] = 0
    result['data'] = parent_dictionary
    return Response(result)


# 递归取字典
def two_level_dictionary(parent_dist):
    objs = data_dictionary.objects.filter(parent_id=parent_dist['code'])
    for obj in objs:
        dist = {}
        dist['id'] = obj.id
        dist['dictionary_name'] = obj.dictionary_name
        dist['parent_id'] = obj.parent_id
        dist['order_num'] = obj.order_num
        dist['code'] = obj.code
        dist['status'] = obj.status
        dist['description'] = obj.description
        dist['nodes'] = []
        parent_dist['nodes'].append(dist)
        two_level_dictionary(dist)
    return parent_dist
# 递归删除字典
def delete_dictionary(parent_id):
    objs = data_dictionary.objects.filter(parent_id=parent_id)
    if objs:
        for obj in objs:
            data_dictionary.objects.filter(code=obj.code).delete()
            delete_dictionary(obj.code)
    else:
        data_dictionary.objects.filter(code=parent_id).delete()





















